from pathlib import Path
import textwrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "assets" / "diagrams"
OUT_FILE = OUT_DIR / "phoenix_architecture_diagrams_editable.pptx"

PX_PER_INCH = 144

COLORS = {
    "bg": RGBColor(11, 17, 30),
    "panel": RGBColor(22, 32, 50),
    "line": RGBColor(83, 104, 132),
    "text": RGBColor(239, 246, 255),
    "muted": RGBColor(166, 180, 202),
    "label_bg": RGBColor(18, 27, 43),
    "label_text": RGBColor(239, 246, 255),
    "dark_text": RGBColor(8, 14, 25),
    "cyan": RGBColor(42, 196, 255),
    "blue": RGBColor(72, 128, 255),
    "green": RGBColor(78, 215, 137),
    "amber": RGBColor(255, 194, 79),
    "red": RGBColor(255, 91, 105),
    "purple": RGBColor(170, 125, 255),
}


def p(value):
    return Inches(value / PX_PER_INCH)


def add_textbox(slide, x, y, w, h, text, size, color, bold=False, name=None):
    shape = slide.shapes.add_textbox(p(x), p(y), p(w), p(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    for idx, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.name = "Aptos"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        para.space_after = Pt(2)
    return shape


def add_base(slide, title, subtitle):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLORS["bg"]
    add_textbox(slide, 80, 58, 1600, 72, title, 33, COLORS["text"], True, "title")
    add_textbox(slide, 84, 132, 1500, 40, subtitle, 16, COLORS["muted"], False, "subtitle")
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        p(80),
        p(178),
        p(300),
        p(10),
    )
    bar.name = "title accent"
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["cyan"]
    bar.line.fill.background()


def wrap_body(text):
    lines = []
    for part in text.split("\n"):
        lines.extend(textwrap.wrap(part, width=31) or [""])
    return "\n".join(lines)


def add_card(slide, x, y, title, body, accent, width=360, height=180, name=None):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        p(x),
        p(y),
        p(width),
        p(height),
    )
    card.name = name or title
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["panel"]
    card.line.color.rgb = COLORS["line"]
    card.line.width = Pt(1.1)

    stripe = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        p(x),
        p(y),
        p(11),
        p(height),
    )
    stripe.name = f"{card.name} accent"
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = accent
    stripe.line.fill.background()

    add_textbox(
        slide,
        x + 30,
        y + 24,
        width - 55,
        42,
        title,
        19,
        COLORS["text"],
        True,
        f"{card.name} title",
    )
    add_textbox(
        slide,
        x + 30,
        y + 78,
        width - 55,
        height - 90,
        wrap_body(body),
        12.5,
        COLORS["muted"],
        False,
        f"{card.name} body",
    )
    return (x, y, x + width, y + height)


def add_label(slide, x, y, text, color, dark=False, name=None):
    width = max(95, len(text) * 8.2 + 32)
    height = 34
    label = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        p(x),
        p(y),
        p(width),
        p(height),
    )
    label.name = name or f"label {text}"
    label.fill.solid()
    label.fill.fore_color.rgb = color if dark else COLORS["label_bg"]
    label.line.color.rgb = color
    label.line.width = Pt(0.8)
    tf = label.text_frame
    tf.clear()
    tf.margin_left = p(10)
    tf.margin_right = p(10)
    tf.margin_top = p(3)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9.8)
    run.font.color.rgb = COLORS["dark_text"] if dark else COLORS["label_text"]
    return label


def add_pill(slide, x, y, text, color, name=None):
    width = len(text) * 7.4 + 36
    pill = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        p(x),
        p(y),
        p(width),
        p(34),
    )
    pill.name = name or f"note {text[:20]}"
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.fill.background()
    tf = pill.text_frame
    tf.clear()
    tf.margin_left = p(16)
    tf.margin_right = p(16)
    tf.margin_top = p(4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(9.8)
    run.font.color.rgb = COLORS["dark_text"]
    return pill


def add_arrow(slide, start, end, color, width=2.8, label=None, label_offset=(0, -38), name=None):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        p(start[0]),
        p(start[1]),
        p(end[0]),
        p(end[1]),
    )
    line.name = name or f"arrow {label or ''}".strip()
    line.line.color.rgb = color
    line.line.width = Pt(width)
    try:
        line.line.end_arrowhead = 3
    except Exception:
        pass
    if label:
        lx = (start[0] + end[0]) / 2 + label_offset[0]
        ly = (start[1] + end[1]) / 2 + label_offset[1]
        add_label(slide, lx, ly, label, color, False, f"{line.name} label")
    return line


def add_dashed_hint(slide, x1, y1, x2, y2, color, name):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, p(x1), p(y1), p(x2), p(y2))
    line.name = name
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return line


def auth_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(
        slide,
        "Splash, Onboarding and Auth Relationship",
        "Code path: GoRouter -> SplashScreen -> AuthCubit -> Onboarding/Auth -> Home",
    )
    b1 = add_card(slide, 110, 330, "Splash Screen", "Initial route /splash\nLogo animation\ncheckAuthStatus()", COLORS["cyan"])
    b2 = add_card(slide, 560, 240, "Auth Gate", "AuthCubit reads local flag\nthen FirebaseAuth currentUser\nand users/{uid}", COLORS["amber"], 420, 220)
    b3 = add_card(slide, 1100, 220, "Home", "AuthSuccess routes to /home\nDroneMainShell opens main app", COLORS["green"], 390, 180)
    b4 = add_card(slide, 560, 610, "Onboarding", "4-page PageView\nSkip or Start now\nroutes to /sign-in", COLORS["purple"], 420, 210)
    b5 = add_card(slide, 1110, 610, "Auth Screens", "Sign in, sign up, OTP\nGoogle sign-in\npassword reset", COLORS["red"], 400, 210)
    b6 = add_card(slide, 1540, 610, "Firebase Auth", "Email/password\nGoogle credential\nFirestore user doc", COLORS["amber"], 300, 210)

    add_arrow(slide, (b1[2], 420), (b2[0], 350), COLORS["cyan"], label="on app start", label_offset=(-70, -72))
    add_arrow(slide, (b2[2], 320), (b3[0], 310), COLORS["green"], label="logged in")
    add_arrow(slide, (770, b2[3]), (770, b4[1]), COLORS["purple"], label="not logged in", label_offset=(-55, -5))
    add_arrow(slide, (b4[2], 715), (b5[0], 715), COLORS["purple"], label="skip / start")
    add_arrow(slide, (b5[2], 715), (b6[0], 715), COLORS["amber"], label="credentials")
    add_arrow(slide, (1310, b5[1]), (1300, b3[3]), COLORS["green"], label="AuthSuccess", label_offset=(-120, -20))
    add_pill(slide, 112, 925, "Routes: /splash -> /onboarding -> /sign-in|/sign-up|/otp|/forgot-password -> /home", COLORS["cyan"])
    add_pill(slide, 112, 970, "Storage: LocalStorageService setLoggedIn(true) after successful auth", COLORS["green"])


def mapping_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(
        slide,
        "Mapping With Firebase and Mission Flow",
        "Code path: DroneRepositoryImpl streams Firestore docs into Cubits and map widgets",
    )
    firebase = add_card(slide, 760, 260, "Cloud Firestore", "drone/location\ndrone/status\ndrone/reports/entries\ndrone/commands", COLORS["amber"], 420, 300)
    pi = add_card(slide, 110, 270, "Raspberry Pi Sync", "firebase_sync.py\nwrites sensors/status\nwrites location every 10s\nreads commands", COLORS["green"], 430, 290)
    repo = add_card(slide, 1350, 230, "Flutter Repository", "watchDroneLocation()\nwatchDroneStatus()\nwatchDroneReports()\nsendCommand()", COLORS["cyan"], 430, 310)
    cubits = add_card(slide, 1120, 650, "Cubits", "DroneTrackingCubit keeps path\nDroneStatusCubit updates stats\nVideoFeedCubit controls mode", COLORS["purple"], 400, 210)
    mapv = add_card(slide, 540, 650, "Map UI", "FlutterMap + OSM tiles\nDrone marker\nUser marker\nPolyline path history", COLORS["blue"], 420, 230)
    mission = add_card(slide, 112, 650, "Mission Button", "StartMissionButton\nsends start_mission\ncommand to Firestore", COLORS["red"], 360, 210)

    add_arrow(slide, (pi[2], 360), (firebase[0], 360), COLORS["green"], label="status + location")
    add_arrow(slide, (firebase[2], 370), (repo[0], 370), COLORS["cyan"], label="snapshots")
    add_arrow(slide, (repo[0] + 180, repo[3]), (cubits[0] + 240, cubits[1]), COLORS["purple"], label="streams")
    add_arrow(slide, (cubits[0], 755), (mapv[2], 755), COLORS["blue"], label="lat/lng + path")
    add_arrow(slide, (mission[2], 755), (firebase[0], 500), COLORS["red"], label="commands")
    add_arrow(slide, (firebase[0], 455), (pi[2], 455), COLORS["amber"], label="poll command")
    add_pill(slide, 80, 930, "Map center defaults to Cairo when Firestore location doc is missing", COLORS["amber"])
    add_pill(slide, 80, 975, "Disconnect state is emitted if no location update arrives for 10 seconds", COLORS["red"])


def system_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_base(
        slide,
        "Raspberry Pi, Video, Laptop Server and Firebase",
        "Code path: camera_stream.py -> laptop_server TCP/FastAPI -> Flutter WebSockets + Firestore telemetry",
    )
    pi = add_card(slide, 80, 260, "Raspberry Pi 5", "Camera module\nDHT11, MQ-2, MQ-8\nMPU6050, IR sensors\nGPIO motors", COLORS["green"], 420, 310)
    tcp = add_card(slide, 610, 230, "Laptop FastAPI Server", "TCP receiver :9000\nOpenCV frame decode\nAI pipeline\nWS endpoints :8000", COLORS["cyan"], 460, 340)
    app = add_card(slide, 1320, 220, "Flutter Mobile App", "WsClient connects to\n/ws/video\n/ws/detections\n/ws/commands", COLORS["blue"], 430, 300)
    fb = add_card(slide, 620, 690, "Firebase", "status + telemetry\nlocation\nreports\ncommands", COLORS["amber"], 420, 240)
    control = add_card(slide, 1320, 690, "Control Paths", "Joystick sends move\nStart mission command\nPi executes GPIO steps", COLORS["red"], 430, 220)

    add_arrow(slide, (pi[2], 345), (tcp[0], 345), COLORS["cyan"], label="JPEG frames over TCP", label_offset=(-72, -58))
    add_arrow(slide, (tcp[2], 325), (app[0], 325), COLORS["blue"], label="video bytes")
    add_arrow(slide, (tcp[2], 455), (app[0], 455), COLORS["purple"], label="detections JSON")
    add_arrow(slide, (app[0] + 40, app[3]), (control[0] + 40, control[1]), COLORS["red"], label="joystick / mission UI", label_offset=(-150, 12))
    add_arrow(slide, (control[0], 805), (fb[2], 805), COLORS["red"], label="Firestore command", label_offset=(-135, -52))
    add_arrow(slide, (fb[0], 810), (pi[2] - 5, 520), COLORS["amber"], label="Pi reads command", label_offset=(-25, -58))
    add_arrow(slide, (pi[0] + 300, pi[3]), (fb[0], 730), COLORS["green"], label="sensor sync", label_offset=(-95, 18))
    add_arrow(slide, (fb[2] - 20, 715), (app[0] + 210, app[3]), COLORS["amber"], label="map/status streams", label_offset=(-95, -54))
    add_pill(slide, 80, 930, "Video is real-time via TCP -> WebSocket, separate from Firebase telemetry", COLORS["cyan"])
    add_pill(slide, 80, 975, "Firebase keeps map, status, reports, and mission commands synchronized", COLORS["amber"])


def build_pptx():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    auth_slide(prs)
    mapping_slide(prs)
    system_slide(prs)
    prs.save(OUT_FILE)
    print(OUT_FILE)


if __name__ == "__main__":
    build_pptx()
