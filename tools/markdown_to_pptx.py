from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]


COLORS = {
    "bg": RGBColor(12, 17, 28),
    "panel": RGBColor(21, 31, 48),
    "accent": RGBColor(42, 196, 255),
    "accent2": RGBColor(255, 87, 87),
    "text": RGBColor(238, 244, 250),
    "muted": RGBColor(162, 176, 194),
    "line": RGBColor(57, 73, 96),
}


def _set_run(run, size=24, bold=False, color=None):
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or COLORS["text"]


def _strip_inline(text):
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = text.replace("**", "")
    text = text.replace("*", "")
    return text.strip()


def _split_marp_slides(markdown):
    parts = re.split(r"^---\s*$", markdown, flags=re.MULTILINE)
    slides = []
    for part in parts:
        part = re.sub(r"<!--.*?-->", "", part, flags=re.DOTALL).strip()
        if not part:
            continue
        if part.startswith("marp:") or "\ntheme:" in part[:120]:
            continue
        slides.append(part)
    return slides


def _parse_slide(block):
    lines = [line.rstrip() for line in block.splitlines()]
    title = "Phoenix"
    subtitle = None
    body = []
    in_code = False
    code_lines = []

    for line in lines:
        if line.startswith("```"):
            in_code = not in_code
            if not in_code and code_lines:
                body.append(("code", "\n".join(code_lines)))
                code_lines = []
            continue
        if in_code:
            code_lines.append(line)
            continue
        if line.startswith("# "):
            title = _strip_inline(line[2:])
        elif line.startswith("## "):
            subtitle = _strip_inline(line[3:])
        elif line.startswith("!["):
            continue
        elif line.strip():
            body.append(("text", line.strip()))
    return title, subtitle, body


def _add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]


def _add_footer(slide, index):
    left = Inches(0.45)
    top = Inches(7.08)
    width = Inches(12.4)
    height = Inches(0.22)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Phoenix | {index}"
    _set_run(run, size=9, color=COLORS["muted"])


def _add_title(slide, title, subtitle=None):
    shape = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12), Inches(0.8))
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=34, bold=True)

    bar = slide.shapes.add_shape(1, Inches(0.68), Inches(1.24), Inches(1.2), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLORS["accent"]
    bar.line.color.rgb = COLORS["accent"]

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.68), Inches(1.36), Inches(11.8), Inches(0.45))
        p = sub.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        _set_run(run, size=18, color=COLORS["muted"])


def _add_body(slide, body, start_top=1.75):
    box = slide.shapes.add_textbox(Inches(0.85), Inches(start_top), Inches(11.8), Inches(5.0))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True

    current = tf.paragraphs[0]
    first = True
    in_table = False

    for kind, raw in body:
        text = raw.strip()
        if not text:
            continue
        if text.startswith("|"):
            if not in_table:
                p = current if first else tf.add_paragraph()
                p.text = "Key technologies and responsibilities are grouped by layer."
                p.level = 0
                p.font.size = Pt(20)
                p.font.color.rgb = COLORS["text"]
                first = False
                in_table = True
            continue
        in_table = False

        if kind == "code":
            p = current if first else tf.add_paragraph()
            p.text = text
            p.font.name = "Cascadia Mono"
            p.font.size = Pt(15)
            p.font.color.rgb = COLORS["accent"]
            first = False
            continue

        clean = _strip_inline(text)
        if clean.startswith("- "):
            p = current if first else tf.add_paragraph()
            p.text = clean[2:]
            p.level = 0
            p.font.size = Pt(22)
            p.font.color.rgb = COLORS["text"]
        elif re.match(r"^\d+\. ", clean):
            p = current if first else tf.add_paragraph()
            p.text = re.sub(r"^\d+\. ", "", clean)
            p.level = 0
            p.font.size = Pt(22)
            p.font.color.rgb = COLORS["text"]
        else:
            p = current if first else tf.add_paragraph()
            p.text = clean
            p.level = 0
            p.font.size = Pt(21)
            p.font.color.rgb = COLORS["muted"]
            p.space_after = Pt(8)
        first = False


def _add_cover(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.6), Inches(1.0))
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=48, bold=True)

    sub = slide.shapes.add_textbox(Inches(0.85), Inches(2.62), Inches(10.8), Inches(0.55))
    p = sub.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    _set_run(run, size=24, color=COLORS["accent"])

    desc = slide.shapes.add_textbox(Inches(0.85), Inches(3.45), Inches(9.8), Inches(0.6))
    p = desc.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Real-time ground control station for AI-assisted drone missions."
    _set_run(run, size=22, color=COLORS["muted"])

    panel = slide.shapes.add_shape(1, Inches(0.85), Inches(4.55), Inches(4.7), Inches(0.08))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["accent2"]
    panel.line.color.rgb = COLORS["accent2"]
    return slide


def build_project_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    markdown = (ROOT / "phoenix_project_presentation.md").read_text(encoding="utf-8")
    slides = _split_marp_slides(markdown)

    for idx, block in enumerate(slides, start=1):
        title, subtitle, body = _parse_slide(block)
        if idx == 1:
            slide = _add_cover(prs, title, subtitle or "Drone Search & Rescue System")
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _add_background(slide)
            _add_title(slide, title, subtitle)
            _add_body(slide, body)
        _add_footer(slide, idx)

    out = ROOT / "Phoenix_Drone_Search_Rescue_Presentation.pptx"
    prs.save(out)
    return out


def _parse_qna(markdown):
    sections = []
    current = None
    for line in markdown.splitlines():
        line = line.strip()
        if line.startswith("## "):
            if current:
                sections.append(current)
            current = [_strip_inline(line[3:]), []]
        elif line and not line.startswith("#") and current:
            current[1].append(_strip_inline(line))
    if current:
        sections.append(current)
    return sections


def build_qna_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    markdown = (ROOT / "phoenix_presentation_qna.md").read_text(encoding="utf-8")
    sections = _parse_qna(markdown)

    slide = _add_cover(prs, "Phoenix Q&A", "Expected discussion questions")
    _add_footer(slide, 1)

    for idx, (question, answer_lines) in enumerate(sections, start=2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_background(slide)
        _add_title(slide, question)
        body = [("text", line) for line in answer_lines]
        _add_body(slide, body, start_top=1.9)
        _add_footer(slide, idx)

    out = ROOT / "Phoenix_Presentation_QA.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    project = build_project_pptx()
    qna = build_qna_pptx()
    print(project)
    print(qna)
