from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "Phoenix.pptx"
LOGO = ROOT / "assets/images/splash/phoenix_3.png"


class C:
    bg = RGBColor(8, 13, 24)
    panel = RGBColor(17, 26, 43)
    panel2 = RGBColor(24, 36, 58)
    cyan = RGBColor(34, 211, 238)
    blue = RGBColor(59, 130, 246)
    green = RGBColor(34, 197, 94)
    orange = RGBColor(251, 146, 60)
    red = RGBColor(248, 82, 82)
    purple = RGBColor(168, 85, 247)
    white = RGBColor(244, 250, 255)
    muted = RGBColor(156, 172, 193)
    dim = RGBColor(83, 101, 128)
    black = RGBColor(5, 8, 14)
    shadow = RGBColor(3, 7, 15)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C.bg
    return s


def text(
    s,
    x,
    y,
    w,
    h,
    value,
    size=20,
    color=C.white,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(s, value, subtitle=None):
    text(s, 0.65, 0.35, 11.6, 0.52, value, 28, C.white, True)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.67), Inches(0.98), Inches(1.0), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C.cyan
    bar.line.fill.background()
    if subtitle:
        text(s, 0.67, 1.08, 11.2, 0.36, subtitle, 14, C.muted)


def footer(s, n):
    text(s, 11.8, 7.05, 0.9, 0.22, f"{len(prs.slides)}/{TOTAL_SLIDES}", 9, C.dim, align=PP_ALIGN.RIGHT)


def soft_panel(s, x, y, w, h, radius=True, fill=C.panel):
    shadow = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x + 0.055),
        Inches(y + 0.065),
        Inches(w),
        Inches(h),
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = C.shadow
    shadow.line.fill.background()

    shape = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def box(s, x, y, w, h, label, sub="", fill=C.panel, accent=C.cyan, size=17):
    shape = soft_panel(s, x, y, w, h, fill=fill)
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.8)

    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.12), Inches(y + 0.09), Inches(w - 0.24), Inches(0.045))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()

    text(
        s,
        x + 0.16,
        y + 0.2,
        w - 0.32,
        0.34,
        label,
        size,
        C.white,
        True,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    if sub:
        text(
            s,
            x + 0.2,
            y + 0.58,
            w - 0.4,
            h - 0.68,
            sub,
            11,
            C.muted,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    return shape


def card(s, x, y, w, h, head, body, accent=C.cyan):
    shape = soft_panel(s, x, y, w, h)
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.1)

    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.12), Inches(y + 0.1), Inches(w - 0.24), Inches(0.055))
    strip.fill.solid()
    strip.fill.fore_color.rgb = accent
    strip.line.fill.background()

    text(s, x + 0.24, y + 0.24, w - 0.48, 0.35, head, 16, C.white, True, valign=MSO_ANCHOR.MIDDLE)
    text(s, x + 0.24, y + 0.7, w - 0.48, h - 0.82, body, 12, C.muted, valign=MSO_ANCHOR.MIDDLE)
    return shape


def arrow(s, x1, y1, x2, y2, color=C.cyan, marker=False):
    line = s.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(2.6)
    if marker:
        end = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x2 - 0.055), Inches(y2 - 0.055), Inches(0.11), Inches(0.11))
        end.fill.solid()
        end.fill.fore_color.rgb = color
        end.line.fill.background()
    return line


def bullets(s, x, y, items, size=17, color=C.white, gap=0.44):
    for i, item in enumerate(items):
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + i * gap + 0.08), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = C.cyan
        dot.line.fill.background()
        text(s, x + 0.22, y + i * gap, 10.8, 0.25, item, size, color)


def label_chip(s, x, y, w, value, color):
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.26))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C.black
    shape.line.color.rgb = color
    shape.line.width = Pt(0.8)
    text(s, x + 0.06, y + 0.02, w - 0.12, 0.2, value, 10, color, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    return shape


def add_logo(s, x, y, size):
    if LOGO.exists():
        s.shapes.add_picture(str(LOGO), Inches(x), Inches(y), Inches(size), Inches(size))


def phone_mock(s, x, y, w, h, mode="dashboard"):
    outer = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    outer.fill.solid()
    outer.fill.fore_color.rgb = C.black
    outer.line.color.rgb = C.dim
    inner = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.13), Inches(y + 0.22), Inches(w - 0.26), Inches(h - 0.44))
    inner.fill.solid()
    inner.fill.fore_color.rgb = RGBColor(12, 20, 34)
    inner.line.fill.background()
    if mode == "splash":
        add_logo(s, x + 0.62, y + 1.0, 1.25)
        text(s, x + 0.28, y + 2.45, w - 0.56, 0.32, "PHOENIX", 17, C.white, True, PP_ALIGN.CENTER)
        text(s, x + 0.28, y + 2.88, w - 0.56, 0.44, "loading app state", 10, C.muted, align=PP_ALIGN.CENTER)
    elif mode == "onboarding":
        card(s, x + 0.28, y + 0.55, w - 0.56, 1.05, "INTRO", "rescue mission value", C.cyan)
        card(s, x + 0.28, y + 1.9, w - 0.56, 1.05, "FEATURES", "video | map | AI", C.orange)
        card(s, x + 0.28, y + 3.25, w - 0.56, 0.65, "GET STARTED", "continue", C.green)
    elif mode == "auth":
        card(s, x + 0.28, y + 0.55, w - 0.56, 0.85, "SIGN IN", "email + password", C.cyan)
        card(s, x + 0.28, y + 1.65, w - 0.56, 0.85, "GOOGLE", "social login", C.orange)
        card(s, x + 0.28, y + 2.75, w - 0.56, 0.85, "RESET", "forgot password", C.purple)
    elif mode == "dashboard":
        card(s, x + 0.28, y + 0.5, w - 0.56, 1.18, "LIVE FEED", "AI Overlay active", C.cyan)
        card(s, x + 0.28, y + 1.9, w - 0.56, 0.85, "STATUS", "Battery 82% | Alt 34m", C.green)
        card(s, x + 0.28, y + 2.95, w - 0.56, 0.85, "REPORTS", "2 humans detected", C.orange)
    elif mode == "map":
        map_mock(s, x + 0.28, y + 0.5, w - 0.56, h - 1.0)
    else:
        video_mock(s, x + 0.28, y + 0.5, w - 0.56, h - 1.0)


def map_mock(s, x, y, w, h):
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(18, 45, 53)
    bg.line.color.rgb = C.cyan
    for i in range(5):
        arrow(s, x + 0.3, y + 0.45 + i * 0.5, x + w - 0.3, y + 0.25 + i * 0.35, C.dim)
    path = s.shapes.add_connector(MSO_CONNECTOR.CURVE, Inches(x + 0.45), Inches(y + h - 0.5), Inches(x + w - 0.5), Inches(y + 0.55))
    path.line.color.rgb = C.orange
    path.line.width = Pt(3)
    drone = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x + w - 0.72), Inches(y + 0.42), Inches(0.28), Inches(0.28))
    drone.fill.solid()
    drone.fill.fore_color.rgb = C.red
    drone.line.fill.background()
    user = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.55), Inches(y + h - 0.85), Inches(0.22), Inches(0.22))
    user.fill.solid()
    user.fill.fore_color.rgb = C.blue
    user.line.color.rgb = C.white


def video_mock(s, x, y, w, h):
    bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(10, 25, 36)
    bg.line.color.rgb = C.cyan
    for i in range(4):
        arrow(s, x + 0.1, y + 0.5 + i * 0.45, x + w - 0.1, y + 0.5 + i * 0.45, C.dim)
    target = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + w * 0.56), Inches(y + h * 0.27), Inches(w * 0.22), Inches(h * 0.48))
    target.fill.background()
    target.line.color.rgb = C.red
    target.line.width = Pt(2.2)
    text(s, x + w * 0.55, y + h * 0.19, w * 0.28, 0.22, "human 92%", 10, C.red, True, PP_ALIGN.CENTER)


def three_layer_diagram(s, y=2.05):
    box(s, 0.8, y, 2.45, 1.15, "Raspberry Pi", "GPS + Camera\nTelemetry", C.panel, C.green)
    box(s, 3.95, y, 2.35, 1.15, "Firebase", "Auth + Firestore\nLive streams", C.panel, C.orange)
    box(s, 6.95, y, 2.65, 1.15, "FastAPI Server", "TCP receiver\nAI pipeline", C.panel, C.cyan)
    box(s, 10.25, y, 2.3, 1.15, "Flutter App", "Video + Map\nControls", C.panel, C.blue)
    arrow(s, 3.25, y + 0.45, 3.95, y + 0.45, C.orange, True)
    arrow(s, 3.25, y + 0.82, 6.95, y + 0.82, C.cyan, True)
    arrow(s, 9.6, y + 0.55, 10.25, y + 0.55, C.cyan, True)
    arrow(s, 10.25, y + 1.05, 6.95, y + 1.05, C.purple, True)


slides = []


def s01():
    s = slide()
    add_logo(s, 8.7, 0.75, 3.4)
    text(s, 0.85, 1.2, 7.6, 0.8, "PHOENIX", 58, C.white, True)
    text(s, 0.9, 2.08, 7.2, 0.42, "Drone Search & Rescue System", 24, C.cyan, True)
    text(s, 0.92, 2.9, 6.6, 0.78, "AI-assisted ground control for live drone missions.", 24, C.muted)
    card(s, 0.92, 4.35, 2.35, 1.05, "Live Video", "Drone camera stream", C.cyan)
    card(s, 3.55, 4.35, 2.35, 1.05, "AI Detection", "Human + thermal alerts", C.red)
    card(s, 6.18, 4.35, 2.35, 1.05, "GPS Map", "Real-time tracking", C.green)
    footer(s, 1)


def s02():
    s = slide()
    title(s, "Problem", "Search & rescue teams need faster field awareness")
    card(s, 0.9, 2.0, 3.45, 2.25, "Dangerous Areas", "Rescuers cannot quickly enter collapsed, flooded, or unstable zones.", C.red)
    card(s, 4.95, 2.0, 3.45, 2.25, "Limited Visibility", "Ground teams miss victims hidden by distance, smoke, or terrain.", C.orange)
    card(s, 9.0, 2.0, 3.45, 2.25, "Slow Decisions", "Manual searching delays mission coordination and victim detection.", C.cyan)
    text(s, 1.4, 5.45, 10.6, 0.5, "Goal: turn drone data into instant, actionable rescue information.", 24, C.white, True, PP_ALIGN.CENTER)
    footer(s, 2)


def s03():
    s = slide()
    title(s, "Solution", "A mobile ground control station for rescue missions")
    three_layer_diagram(s, 2.15)
    bullets(s, 1.2, 4.55, ["One app for video, map, telemetry, reports, and control", "Laptop server handles streaming and AI processing", "Firebase keeps GPS/status updates live and reliable"], 18)
    footer(s, 3)


def s03a():
    s = slide()
    title(s, "App Entry Point", "Splash screen and startup decision")
    phone_mock(s, 0.95, 1.55, 2.65, 5.0, "splash")
    box(s, 4.25, 1.9, 2.25, 0.85, "Initialize", "Firebase + services", accent=C.cyan)
    box(s, 7.0, 1.9, 2.25, 0.85, "Check Session", "Firebase user + local flag", accent=C.orange)
    box(s, 9.75, 1.9, 2.25, 0.85, "Route", "onboarding/auth/home", accent=C.green)
    arrow(s, 6.5, 2.32, 7.0, 2.32, C.orange, True)
    arrow(s, 9.25, 2.32, 9.75, 2.32, C.green, True)
    bullets(s, 4.35, 4.35, ["First screen controls the user journey", "Decides whether to continue to onboarding, authentication, or home"], 18)
    footer(s, 0)


def s03b():
    s = slide()
    title(s, "Onboarding Flow", "Introduce the rescue system before login")
    phone_mock(s, 0.95, 1.55, 2.65, 5.0, "onboarding")
    card(s, 4.25, 1.75, 3.1, 1.0, "Purpose", "explain search & rescue value", C.cyan)
    card(s, 8.15, 1.75, 3.1, 1.0, "Capabilities", "video, map, AI, control", C.orange)
    card(s, 4.25, 3.35, 3.1, 1.0, "Page Indicator", "multi-page intro", C.purple)
    card(s, 8.15, 3.35, 3.1, 1.0, "Next Step", "send user to sign in", C.green)
    footer(s, 0)


def s03c():
    s = slide()
    title(s, "Authentication", "Secure access before mission control")
    phone_mock(s, 0.95, 1.55, 2.65, 5.0, "auth")
    box(s, 4.25, 1.9, 2.4, 0.9, "Sign Up", "email verification", accent=C.cyan)
    box(s, 7.25, 1.9, 2.4, 0.9, "Sign In", "Firebase Auth", accent=C.green)
    box(s, 10.25, 1.9, 2.4, 0.9, "Recovery", "reset password", accent=C.purple)
    card(s, 4.25, 4.15, 3.3, 1.15, "User Document", "created in Firestore after auth success", C.orange)
    card(s, 8.35, 4.15, 3.3, 1.15, "Local Session", "shared preferences login state", C.blue)
    footer(s, 0)


def s03d():
    s = slide()
    title(s, "Navigation Flow", "Routes used by the Flutter app")
    steps = [
        ("Splash", "/splash", C.cyan),
        ("Onboarding", "/onboarding", C.orange),
        ("Auth", "/sign-in", C.green),
        ("Home", "/home", C.purple),
        ("Video", "/fullscreen-video", C.red),
    ]
    for i, (head, route, color) in enumerate(steps):
        x = 0.95 + i * 2.42
        box(s, x, 2.45, 1.78, 0.95, head, route, accent=color)
        if i < len(steps) - 1:
            arrow(s, x + 1.78, 2.92, x + 2.42, 2.92, C.dim, True)
    bullets(s, 1.15, 4.95, ["go_router keeps screens declarative and easy to follow", "Main mission screens are reached only after entry/auth flow"], 18)
    footer(s, 0)


def s04():
    s = slide()
    title(s, "System Block Diagram", "Main components and communication paths")
    three_layer_diagram(s, 1.85)
    label_chip(s, 1.18, 3.62, 1.65, "GPS / Status", C.orange)
    label_chip(s, 4.1, 3.62, 1.85, "Firestore stream", C.orange)
    label_chip(s, 5.32, 3.94, 1.35, "TCP video", C.cyan)
    label_chip(s, 8.95, 3.62, 1.45, "WebSocket", C.cyan)
    label_chip(s, 8.52, 3.94, 1.35, "Commands", C.purple)
    footer(s, 4)


def s05():
    s = slide()
    title(s, "Data Flow", "Separate paths for speed and reliability")
    box(s, 0.8, 2.0, 2.0, 0.8, "GPS", "lat/lng/status", accent=C.orange)
    box(s, 3.35, 2.0, 2.1, 0.8, "Firebase", "snapshots", accent=C.orange)
    box(s, 6.0, 2.0, 2.2, 0.8, "Flutter Map", "marker + path", accent=C.green)
    arrow(s, 2.8, 2.4, 3.35, 2.4, C.orange, True)
    arrow(s, 5.45, 2.4, 6.0, 2.4, C.orange, True)
    box(s, 0.8, 3.65, 2.0, 0.8, "Camera", "JPEG frames", accent=C.cyan)
    box(s, 3.35, 3.65, 2.1, 0.8, "FastAPI", "AI process", accent=C.cyan)
    box(s, 6.0, 3.65, 2.2, 0.8, "Flutter Video", "live feed", accent=C.cyan)
    arrow(s, 2.8, 4.05, 3.35, 4.05, C.cyan, True)
    arrow(s, 5.45, 4.05, 6.0, 4.05, C.cyan, True)
    box(s, 10.0, 2.8, 2.3, 0.95, "Commands", "joystick -> server", accent=C.purple)
    arrow(s, 10.0, 3.25, 8.2, 4.0, C.purple, True)
    footer(s, 5)


def s06():
    s = slide()
    title(s, "Hardware Layer", "Drone onboard module")
    box(s, 1.0, 2.05, 2.25, 1.1, "GPS Module", "live location", accent=C.green)
    box(s, 4.0, 2.05, 2.25, 1.1, "Camera", "video frames", accent=C.cyan)
    box(s, 7.0, 2.05, 2.25, 1.1, "Sensors", "battery/temp", accent=C.orange)
    box(s, 10.0, 2.05, 2.25, 1.1, "Raspberry Pi", "data relay", accent=C.purple)
    arrow(s, 3.25, 2.6, 4.0, 2.6, C.dim, True)
    arrow(s, 6.25, 2.6, 7.0, 2.6, C.dim, True)
    arrow(s, 9.25, 2.6, 10.0, 2.6, C.dim, True)
    bullets(s, 1.15, 4.45, ["Writes location/status to Firebase", "Streams camera frames to laptop server", "Ready for thermal camera integration"], 19)
    footer(s, 6)


def s07():
    s = slide()
    title(s, "Backend Server", "FastAPI middleware on laptop")
    box(s, 0.8, 2.0, 2.0, 0.9, "TCP Receiver", "frames in", accent=C.cyan)
    box(s, 3.25, 2.0, 2.0, 0.9, "Queue", "async buffer", accent=C.blue)
    box(s, 5.7, 2.0, 2.0, 0.9, "AI Pipeline", "detect + draw", accent=C.red)
    box(s, 8.15, 2.0, 2.0, 0.9, "Encoder", "JPEG bytes", accent=C.orange)
    box(s, 10.6, 2.0, 2.0, 0.9, "WebSocket", "broadcast", accent=C.green)
    for x in [2.8, 5.25, 7.7, 10.15]:
        arrow(s, x, 2.45, x + 0.45, 2.45, C.cyan, True)
    card(s, 2.1, 4.25, 2.65, 1.25, "/ws/video", "Processed frames", C.cyan)
    card(s, 5.35, 4.25, 2.65, 1.25, "/ws/detections", "JSON boxes", C.red)
    card(s, 8.6, 4.25, 2.65, 1.25, "/ws/commands", "Control input", C.purple)
    footer(s, 7)


def s08():
    s = slide()
    title(s, "AI Detection", "Human detection and thermal-ready pipeline")
    video_mock(s, 0.9, 1.8, 5.6, 3.4)
    box(s, 7.1, 1.9, 2.05, 0.78, "Frame", "camera input", accent=C.cyan)
    box(s, 9.65, 1.9, 2.05, 0.78, "YOLO", "human boxes", accent=C.red)
    box(s, 7.1, 3.25, 2.05, 0.78, "Thermal", "hot spots", accent=C.orange)
    box(s, 9.65, 3.25, 2.05, 0.78, "Overlay", "HUD output", accent=C.green)
    arrow(s, 9.15, 2.3, 9.65, 2.3, C.red, True)
    arrow(s, 9.15, 3.64, 9.65, 3.64, C.orange, True)
    bullets(s, 7.15, 5.05, ["Dummy fallback keeps demos running", "Final model can replace detector module"], 17)
    footer(s, 8)


def s09():
    s = slide()
    title(s, "Mobile App", "Flutter ground control interface")
    phone_mock(s, 1.0, 1.55, 2.55, 4.95, "dashboard")
    phone_mock(s, 5.35, 1.55, 2.55, 4.95, "map")
    phone_mock(s, 9.7, 1.55, 2.55, 4.95, "video")
    text(s, 1.18, 6.65, 2.2, 0.25, "Dashboard", 14, C.muted, True, PP_ALIGN.CENTER)
    text(s, 5.53, 6.65, 2.2, 0.25, "Live Map", 14, C.muted, True, PP_ALIGN.CENTER)
    text(s, 9.88, 6.65, 2.2, 0.25, "Full Video", 14, C.muted, True, PP_ALIGN.CENTER)
    footer(s, 9)


def s10():
    s = slide()
    title(s, "App Architecture", "Feature-first clean structure")
    box(s, 0.95, 2.0, 2.25, 1.0, "Core", "DI | Router\nTheme | WS", accent=C.cyan)
    box(s, 3.9, 2.0, 2.25, 1.0, "Auth", "Firebase Auth\nGoogle sign-in", accent=C.orange)
    box(s, 6.85, 2.0, 2.25, 1.0, "Drone", "Map | Video\nTelemetry", accent=C.green)
    box(s, 9.8, 2.0, 2.25, 1.0, "Onboarding", "Splash\nIntro flow", accent=C.purple)
    card(s, 2.2, 4.25, 2.6, 1.1, "Domain", "Models + repository contracts", C.cyan)
    card(s, 5.35, 4.25, 2.6, 1.1, "Data", "Firebase + WebSocket implementations", C.orange)
    card(s, 8.5, 4.25, 2.6, 1.1, "Presentation", "Screens + Cubits + widgets", C.green)
    footer(s, 10)


def s11():
    s = slide()
    title(s, "Live Map", "GPS tracking with path history")
    map_mock(s, 0.95, 1.65, 7.0, 4.75)
    bullets(s, 8.55, 2.05, ["Firestore location stream", "OpenStreetMap tiles", "Path history", "Disconnect handling", "Last known position"], 18)
    footer(s, 11)


def s11a():
    s = slide()
    title(s, "Map Data Source", "How GPS reaches the user interface")
    box(s, 0.85, 2.0, 2.1, 0.9, "Raspberry Pi", "reads GPS", accent=C.green)
    box(s, 3.65, 2.0, 2.1, 0.9, "Firestore", "drone/location", accent=C.orange)
    box(s, 6.45, 2.0, 2.1, 0.9, "Repository", "watch stream", accent=C.cyan)
    box(s, 9.25, 2.0, 2.1, 0.9, "Cubit", "emit state", accent=C.purple)
    arrow(s, 2.95, 2.45, 3.65, 2.45, C.orange, True)
    arrow(s, 5.75, 2.45, 6.45, 2.45, C.cyan, True)
    arrow(s, 8.55, 2.45, 9.25, 2.45, C.purple, True)
    card(s, 1.05, 4.25, 3.0, 1.15, "Document", "lat | lng | timestamp", C.orange)
    card(s, 5.15, 4.25, 3.0, 1.15, "Model", "DroneLocation", C.cyan)
    card(s, 9.25, 4.25, 3.0, 1.15, "UI", "marker + route update", C.green)
    footer(s, 0)


def s11b():
    s = slide()
    title(s, "Map Rendering Layers", "What appears on the map")
    map_mock(s, 0.85, 1.65, 6.6, 4.85)
    card(s, 8.05, 1.7, 3.25, 0.88, "Tile Layer", "OpenStreetMap base map", C.cyan)
    card(s, 8.05, 2.85, 3.25, 0.88, "Polyline Layer", "drone path history", C.orange)
    card(s, 8.05, 4.0, 3.25, 0.88, "Marker Layer", "drone + user location", C.red)
    card(s, 8.05, 5.15, 3.25, 0.88, "Control Layer", "center on user button", C.green)
    footer(s, 0)


def s11c():
    s = slide()
    title(s, "Tracking State Machine", "How the app reacts to GPS stream changes")
    box(s, 0.85, 2.25, 2.1, 0.85, "Initial", "screen starts", accent=C.dim)
    box(s, 3.45, 2.25, 2.1, 0.85, "Loading", "subscribe stream", accent=C.cyan)
    box(s, 6.05, 2.25, 2.1, 0.85, "Active", "location + path", accent=C.green)
    box(s, 8.65, 2.25, 2.1, 0.85, "Disconnected", "last known point", accent=C.red)
    arrow(s, 2.95, 2.68, 3.45, 2.68, C.cyan, True)
    arrow(s, 5.55, 2.68, 6.05, 2.68, C.green, True)
    arrow(s, 8.15, 2.68, 8.65, 2.68, C.red, True)
    label_chip(s, 6.2, 3.55, 1.85, "10 sec timeout", C.red)
    bullets(s, 1.1, 4.65, ["Active state updates the map marker instantly", "Path history grows with every new location", "Disconnected state keeps the last known position visible"], 18)
    footer(s, 0)


def s11d():
    s = slide()
    title(s, "Map Screen UX", "Operational view for rescue teams")
    phone_mock(s, 0.95, 1.45, 2.7, 5.15, "map")
    card(s, 4.35, 1.75, 3.15, 1.0, "Header", "mission map context", C.cyan)
    card(s, 8.15, 1.75, 3.15, 1.0, "Status Bar", "battery | temp | signal", C.green)
    card(s, 4.35, 3.3, 3.15, 1.0, "Mission Card", "current operation status", C.orange)
    card(s, 8.15, 3.3, 3.15, 1.0, "Start Button", "send start_mission", C.purple)
    bullets(s, 4.45, 5.35, ["Map is not isolated; it is part of the mission workflow", "Telemetry stays visible while tracking location"], 17)
    footer(s, 0)


def s12():
    s = slide()
    title(s, "Video Feed", "Normal, Thermal, AI Overlay")
    video_mock(s, 0.95, 1.65, 7.1, 4.75)
    card(s, 8.65, 1.85, 2.8, 0.92, "Normal", "raw camera view", C.cyan)
    card(s, 8.65, 3.05, 2.8, 0.92, "Thermal", "heat-focused mode", C.orange)
    card(s, 8.65, 4.25, 2.8, 0.92, "AI Overlay", "boxes + confidence", C.red)
    footer(s, 12)


def s13():
    s = slide()
    title(s, "Drone Control", "Virtual joystick command flow")
    phone_mock(s, 0.95, 1.55, 2.65, 5.0, "video")
    box(s, 4.35, 2.1, 2.1, 0.85, "Joystick", "x/y offset", accent=C.purple)
    box(s, 7.1, 2.1, 2.1, 0.85, "WebSocket", "/ws/commands", accent=C.cyan)
    box(s, 9.85, 2.1, 2.1, 0.85, "Server", "handle command", accent=C.green)
    arrow(s, 6.45, 2.52, 7.1, 2.52, C.purple, True)
    arrow(s, 9.2, 2.52, 9.85, 2.52, C.purple, True)
    bullets(s, 4.45, 4.25, ["move: x/y direction", "neutral command on release", "ready for land / return / stop"], 18)
    footer(s, 13)


def s14():
    s = slide()
    title(s, "Firebase", "Cloud services used by the app")
    card(s, 1.0, 1.8, 3.2, 1.45, "Authentication", "Email/password\nGoogle sign-in", C.orange)
    card(s, 5.05, 1.8, 3.2, 1.45, "Firestore", "Location\nStatus\nReports", C.green)
    card(s, 9.1, 1.8, 3.2, 1.45, "User Data", "Profile document\nLogin state", C.cyan)
    bullets(s, 1.25, 4.6, ["Realtime snapshots update the UI automatically", "Lightweight data stays separate from heavy video streaming"], 20)
    footer(s, 14)


def s14a():
    s = slide()
    title(s, "Telemetry Monitoring", "Live status data shown beside the mission")
    card(s, 0.9, 1.85, 2.65, 1.2, "Battery", "percentage for flight safety", C.green)
    card(s, 3.95, 1.85, 2.65, 1.2, "Altitude", "current drone height", C.cyan)
    card(s, 7.0, 1.85, 2.65, 1.2, "Speed", "movement awareness", C.blue)
    card(s, 10.05, 1.85, 2.65, 1.2, "Temperature", "system/thermal signal", C.orange)
    box(s, 2.05, 4.25, 2.35, 0.9, "Firestore", "drone/status", accent=C.orange)
    box(s, 5.25, 4.25, 2.35, 0.9, "Status Cubit", "stream state", accent=C.purple)
    box(s, 8.45, 4.25, 2.35, 0.9, "UI Widgets", "cards + bars", accent=C.green)
    arrow(s, 4.4, 4.7, 5.25, 4.7, C.purple, True)
    arrow(s, 7.6, 4.7, 8.45, 4.7, C.green, True)
    footer(s, 0)


def s14b():
    s = slide()
    title(s, "Reports & Alerts", "Incident history for mission awareness")
    card(s, 0.95, 1.8, 3.25, 1.25, "Human Detected", "AI or thermal event", C.red)
    card(s, 5.05, 1.8, 3.25, 1.25, "System Overheated", "temperature warning", C.orange)
    card(s, 9.15, 1.8, 3.25, 1.25, "Mission Complete", "operation event", C.green)
    box(s, 1.25, 4.25, 2.6, 0.95, "Firestore", "reports/entries", accent=C.orange)
    box(s, 5.35, 4.25, 2.6, 0.95, "Reports Section", "latest 20 entries", accent=C.cyan)
    box(s, 9.45, 4.25, 2.6, 0.95, "Rescue Team", "quick review", accent=C.green)
    arrow(s, 3.85, 4.72, 5.35, 4.72, C.cyan, True)
    arrow(s, 7.95, 4.72, 9.45, 4.72, C.green, True)
    footer(s, 0)


def s14c():
    s = slide()
    title(s, "Mission Workflow", "From planning to active operation")
    steps = [
        ("Plan", "choose target area", C.cyan),
        ("Start", "send start_mission", C.green),
        ("Monitor", "map + telemetry", C.orange),
        ("Detect", "AI alerts", C.red),
        ("Control", "joystick commands", C.purple),
        ("Review", "reports log", C.blue),
    ]
    for i, (head, body, color) in enumerate(steps):
        x = 0.75 + i * 2.05
        box(s, x, 2.35, 1.55, 0.95, head, body, accent=color)
        if i < len(steps) - 1:
            arrow(s, x + 1.55, 2.82, x + 2.05, 2.82, C.dim, True)
    bullets(s, 1.1, 4.85, ["The current app already sends start_mission", "Mission planner can later add routes, zones, and waypoints"], 18)
    footer(s, 0)


def s15():
    s = slide()
    title(s, "Technology Stack", "Tools selected for real-time operation")
    stacks = [
        ("Flutter", "mobile UI", C.cyan),
        ("BLoC/Cubit", "state", C.blue),
        ("FastAPI", "server", C.green),
        ("Firebase", "cloud", C.orange),
        ("WebSocket", "live stream", C.purple),
        ("OpenCV", "video AI", C.red),
    ]
    for i, (name, sub, color) in enumerate(stacks):
        x = 1.0 + (i % 3) * 4.05
        y = 1.85 + (i // 3) * 1.8
        card(s, x, y, 3.15, 1.15, name, sub, color)
    footer(s, 15)


def s16():
    s = slide()
    title(s, "Completed Work", "Current implementation status")
    bullets(s, 1.1, 1.75, [
        "Flutter app UI and navigation",
        "Authentication flow",
        "Firestore location/status/reports streams",
        "Live map and telemetry dashboard",
        "FastAPI server with TCP receiver",
        "WebSocket video/detection/command channels",
        "AI overlay integration path",
        "Virtual joystick command sending",
    ], 18, C.white, 0.55)
    footer(s, 16)


def s16a():
    s = slide()
    title(s, "Integration Readiness", "How the full system is tested together")
    box(s, 0.85, 1.9, 2.2, 0.85, "1. Server", "run FastAPI", accent=C.green)
    box(s, 3.45, 1.9, 2.2, 0.85, "2. Drone/Pi", "send GPS + frames", accent=C.orange)
    box(s, 6.05, 1.9, 2.2, 0.85, "3. Mobile", "connect WebSocket", accent=C.cyan)
    box(s, 8.65, 1.9, 2.2, 0.85, "4. Observe", "map + video", accent=C.purple)
    arrow(s, 3.05, 2.32, 3.45, 2.32, C.dim, True)
    arrow(s, 5.65, 2.32, 6.05, 2.32, C.dim, True)
    arrow(s, 8.25, 2.32, 8.65, 2.32, C.dim, True)
    card(s, 1.05, 4.1, 2.8, 1.2, "Health Check", "/health endpoint", C.green)
    card(s, 4.25, 4.1, 2.8, 1.2, "Same Network", "phone + laptop + Pi", C.cyan)
    card(s, 7.45, 4.1, 2.8, 1.2, "Fallback Mode", "dummy AI keeps flow alive", C.red)
    footer(s, 0)


def s17():
    s = slide()
    title(s, "Demo Flow", "What the committee will see")
    steps = ["Sign in", "Open dashboard", "View live feed", "Track on map", "Show AI overlay", "Send joystick command"]
    for i, step in enumerate(steps):
        x = 0.8 + i * 2.05
        box(s, x, 2.55, 1.55, 0.9, f"{i + 1}", step, accent=[C.cyan, C.green, C.orange, C.red, C.purple, C.blue][i])
        if i < len(steps) - 1:
            arrow(s, x + 1.55, 3.0, x + 2.05, 3.0, C.dim, True)
    footer(s, 17)


def s18():
    s = slide()
    title(s, "Feature Challenges", "Problems we faced and how the code solves them")
    card(s, 0.95, 1.65, 3.5, 1.45, "Splash/Auth", "Challenge: startup routing.\nSolution: Firebase + local session check.", C.orange)
    card(s, 4.9, 1.65, 3.5, 1.45, "Map Tracking", "Challenge: live GPS dropouts.\nSolution: stream state + last known position.", C.green)
    card(s, 8.85, 1.65, 3.5, 1.45, "Video Latency", "Challenge: heavy frames.\nSolution: TCP/WS path separate from Firebase.", C.cyan)
    card(s, 0.95, 4.15, 3.5, 1.45, "AI Detection", "Challenge: model availability.\nSolution: detector module with fallback.", C.red)
    card(s, 4.9, 4.15, 3.5, 1.45, "Drone Control", "Challenge: continuous joystick input.\nSolution: x/y commands and neutral on release.", C.purple)
    card(s, 8.85, 4.15, 3.5, 1.45, "Team Integration", "Challenge: app/server/Pi sync.\nSolution: clean Firestore + WebSocket contracts.", C.blue)
    footer(s, 18)


def s19():
    s = slide()
    title(s, "Limitations", "Honest current scope")
    bullets(s, 1.15, 2.0, [
        "Final YOLO model still needs replacement of dummy fallback",
        "Full hardware test depends on Raspberry Pi scripts",
        "Mission planning workflow is not complete yet",
        "Drone command acknowledgement is a next integration step",
    ], 20, C.white, 0.62)
    footer(s, 19)


def s20():
    s = slide()
    title(s, "Future Work", "Next improvements")
    card(s, 0.9, 1.8, 2.8, 1.15, "Final AI", "YOLO + thermal model", C.red)
    card(s, 4.0, 1.8, 2.8, 1.15, "Mission Planner", "routes + zones", C.cyan)
    card(s, 7.1, 1.8, 2.8, 1.15, "Recording", "video + logs", C.orange)
    card(s, 10.2, 1.8, 2.8, 1.15, "Team Mode", "multi-user access", C.green)
    card(s, 2.45, 4.1, 2.8, 1.15, "Reconnect", "offline handling", C.purple)
    card(s, 5.55, 4.1, 2.8, 1.15, "Ack System", "command feedback", C.blue)
    card(s, 8.65, 4.1, 2.8, 1.15, "Hardware", "real drone testing", C.cyan)
    footer(s, 20)


def s21():
    s = slide()
    title(s, "Value", "Why Phoenix matters")
    text(s, 1.15, 2.0, 11.0, 0.65, "Faster awareness. Safer decisions. Smarter rescue missions.", 34, C.white, True, PP_ALIGN.CENTER)
    card(s, 1.2, 4.0, 3.25, 1.25, "Field Ready", "mobile-first dark interface", C.cyan)
    card(s, 5.0, 4.0, 3.25, 1.25, "AI Assisted", "human/thermal detection path", C.red)
    card(s, 8.8, 4.0, 3.25, 1.25, "Scalable", "clean three-tier architecture", C.green)
    footer(s, 21)


def s22():
    s = slide()
    add_logo(s, 5.25, 0.75, 2.8)
    text(s, 0.9, 3.55, 11.55, 0.65, "THANK YOU", 46, C.white, True, PP_ALIGN.CENTER)
    text(s, 1.35, 4.35, 10.65, 0.45, "Questions?", 26, C.cyan, True, PP_ALIGN.CENTER)
    text(s, 1.6, 5.3, 10.1, 0.35, "Phoenix Drone Search & Rescue System", 17, C.muted, align=PP_ALIGN.CENTER)
    footer(s, 22)


BUILDERS = [
    s01,
    s02,
    s03,
    s03a,
    s03b,
    s03c,
    s09,
    s10,
    s04,
    s05,
    s11,
    s11a,
    s11b,
    s11c,
    s11d,
    s14a,
    s14b,
    s14c,
    s12,
    s13,
    s06,
    s07,
    s08,
    s18,
    s22,
]
TOTAL_SLIDES = len(BUILDERS)

for builder in BUILDERS:
    builder()

prs.save(OUT)
print(OUT)
