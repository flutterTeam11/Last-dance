#!/usr/bin/env python3
"""
Phoenix - Drone Search & Rescue System
Professional Graduation Project Presentation Generator
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette (Dark theme, matching app) ──
BG_DARK      = RGBColor(0x0D, 0x0D, 0x1A)
BG_CARD      = RGBColor(0x16, 0x16, 0x2E)
BG_SECTION   = RGBColor(0x1A, 0x1A, 0x35)
ACCENT_CYAN  = RGBColor(0x00, 0xE5, 0xFF)
ACCENT_BLUE  = RGBColor(0x00, 0x96, 0xFF)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB0, 0xC0)
DARK_GRAY    = RGBColor(0x60, 0x60, 0x80)
GRADIENT_END = RGBColor(0x0A, 0x0A, 0x1A)
GREEN        = RGBColor(0x00, 0xE5, 0x80)
ORANGE       = RGBColor(0xFF, 0x8C, 0x00)
RED          = RGBColor(0xFF, 0x33, 0x55)

def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, opacity=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_accent_bar(slide, left, top, width=Inches(0.08), height=Inches(0.5), color=ACCENT_CYAN):
    return add_shape(slide, left, top, width, height, color)

def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT_CYAN, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        # Add bullet via prefix
    return txBox

def add_section_title(slide, title, subtitle=None):
    # Accent bar on the left
    add_accent_bar(slide, Inches(0.6), Inches(0.6), Inches(0.08), Inches(0.7), ACCENT_CYAN)
    add_textbox(slide, Inches(0.9), Inches(0.55), Inches(11), Inches(0.8), title, 
                font_size=36, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.9), Inches(1.2), Inches(11), Inches(0.5), subtitle,
                    font_size=18, color=DARK_GRAY)

def add_icon_card(slide, left, top, width, height, icon_text, title, description, accent_color=ACCENT_CYAN):
    card = add_shape(slide, left, top, width, height, BG_CARD)
    card.shadow.inherit = False
    
    # Top accent line
    add_shape(slide, left, top, width, Inches(0.04), accent_color)
    
    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
        left + Inches(0.3), top + Inches(0.3), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(18)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    
    add_textbox(slide, left + Inches(1.1), top + Inches(0.25), width - Inches(1.4), Inches(0.4),
                title, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(1.1), top + Inches(0.65), width - Inches(1.4), Inches(1.0),
                description, font_size=12, color=LIGHT_GRAY)

def new_slide():
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide)
    return slide

def add_page_number(slide, num, total):
    add_textbox(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.RIGHT)

def add_gradient_line(slide, left, top, width):
    for i in range(int(width / Inches(0.02))):
        ratio = i / (width / Inches(0.02))
        r = int(0x0D + (0x00 - 0x0D) * ratio)
        g = int(0x0D + (0xE5 - 0x0D) * ratio)
        b = int(0x1A + (0xFF - 0x1A) * ratio)
        c = RGBColor(r, g, b)
        add_shape(slide, left + Inches(i * 0.02), top, Inches(0.02), Inches(0.04), c)

TOTAL_SLIDES = 24

# ═══════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════
slide = new_slide()
# Background gradient overlay
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), BG_DARK)
# Decorative elements - large circle
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(-1.5), Inches(7), Inches(7))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x00, 0x65, 0xFF)
circle.line.fill.background()
# Make it semi-transparent by setting alpha
# python-pptx doesn't support alpha directly, so we use a workaround
circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(3), Inches(5), Inches(5))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x7C, 0x3A, 0xED)
circle2.line.fill.background()

# Small decorative dots
for x, y in [(1.5, 1.5), (2.0, 6.5), (11.5, 1.0), (0.5, 4.0)]:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.08), Inches(0.08))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_CYAN
    dot.line.fill.background()

# Title content
add_textbox(slide, Inches(1.0), Inches(1.5), Inches(7), Inches(1.0),
            "PHOENIX", font_size=72, color=WHITE, bold=True)
add_textbox(slide, Inches(1.0), Inches(2.4), Inches(7), Inches(0.6),
            "Drone Search & Rescue System", font_size=36, color=ACCENT_CYAN, bold=True)
add_gradient_line(slide, Inches(1.0), Inches(3.2), Inches(3))
add_textbox(slide, Inches(1.0), Inches(3.5), Inches(7), Inches(0.5),
            "Graduation Project  |  Computer Engineering", font_size=18, color=LIGHT_GRAY)
add_textbox(slide, Inches(1.0), Inches(4.5), Inches(7), Inches(1.5),
            "A three-tier ground control system integrating drone hardware,\n"
            "AI-powered human detection, and real-time mobile monitoring\n"
            "for autonomous Search & Rescue operations.",
            font_size=16, color=DARK_GRAY)

# Bottom info
add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.02), ACCENT_CYAN)
add_textbox(slide, Inches(1.0), Inches(6.9), Inches(5), Inches(0.4),
            "Mahmoud Dahy  |  Supervised by: Dr. ...", font_size=12, color=DARK_GRAY)

# ═══════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Agenda", "Presentation Outline")

items = [
    "01    Problem Statement & Motivation",
    "02    Project Objectives",
    "03    System Architecture Overview",
    "04    Hardware Components (Raspberry Pi Module)",
    "05    Backend Server (FastAPI & AI Pipeline)",
    "06    AI Detection Engine (YOLO Integration)",
    "07    Mobile App – Ground Control Interface",
    "08    Key Features: Video, Map, Telemetry, Joystick",
    "09    Technology Stack",
    "10    Challenges & Solutions",
    "11    Future Enhancements",
    "12    Demonstration",
]
add_bullet_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(5.5),
                   items, font_size=18, color=LIGHT_GRAY, spacing=Pt(10))
add_page_number(slide, 2, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 3: PROBLEM STATEMENT
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Problem Statement", "The Challenge in Search & Rescue Operations")

problems = [
    "Critical delays in locating victims during natural disasters, missing person cases, and emergencies",
    "Traditional SAR teams face limited visibility, dangerous terrain, and slow manual search patterns",
    "Lack of real-time aerial intelligence and automated detection capabilities in the field",
    "High cost & complexity of professional drone systems with limited accessibility for first responders",
]
add_bullet_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(3.5),
                   problems, font_size=16, color=LIGHT_GRAY, spacing=Pt(10))

# Stats cards
stats = [
    ("72%", "of SAR operations\nexceed critical\ntime windows", RED),
    ("60%", "reduction in search\narea coverage\nwith aerial view", ACCENT_CYAN),
    ("3x", "faster victim\nlocalization with\nAI detection", GREEN),
]
for i, (num, desc, color) in enumerate(stats):
    left = Inches(1.5 + i * 3.8)
    card = add_shape(slide, left, Inches(5.0), Inches(3.2), Inches(2.0), BG_CARD)
    add_shape(slide, left, Inches(5.0), Inches(3.2), Inches(0.04), color)
    add_textbox(slide, left + Inches(0.3), Inches(5.2), Inches(2.6), Inches(0.6),
                num, font_size=36, color=color, bold=True)
    add_textbox(slide, left + Inches(0.3), Inches(5.8), Inches(2.6), Inches(0.8),
                desc, font_size=12, color=LIGHT_GRAY)

add_page_number(slide, 3, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 4: OBJECTIVES
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Project Objectives", "What We Aim to Achieve")

objectives = [
    ("01", "Real-Time Aerial Surveillance",
     "Live video streaming from drone camera to mobile ground station with minimal latency"),
    ("02", "AI-Powered Human Detection",
     "Automated identification of persons in distress using YOLO deep learning model"),
    ("03", "Live GPS Tracking & Navigation",
     "Real-time drone position tracking on interactive maps with path history"),
    ("04", "Comprehensive Telemetry Monitoring",
     "Battery, altitude, speed, temperature, and system health at a glance"),
    ("05", "Remote Drone Control",
     "Full touch-based virtual joystick for directional drone commands"),
    ("06", "Mission Management",
     "End-to-end mission lifecycle: planning, execution, incident reporting"),
]

for i, (num, title, desc) in enumerate(objectives):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.0)
    top = Inches(1.8 + row * 2.6)
    add_icon_card(slide, left, top, Inches(3.6), Inches(2.2),
                  num, title, desc,
                  [ACCENT_CYAN, ACCENT_BLUE, ACCENT_PURPLE, GREEN, ORANGE, RED][i])

add_page_number(slide, 4, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 5: SYSTEM ARCHITECTURE OVERVIEW
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "System Architecture", "Three-Tier Distributed System")

# Three columns for three tiers
tiers = [
    ("TIER 1", "DRONE HARDWARE", "Raspberry Pi\nCamera Module\nGPS (NEO-6M)\n\nCaptures video &\nlocation data.\nStreams via TCP &\nFirebase sync.", ACCENT_CYAN),
    ("TIER 2", "LAPTOP SERVER", "FastAPI + Uvicorn\nYOLO AI Detection\nOpenCV Processing\nWebSocket Server\n\nReceives drone feed,\nruns AI inference,\nbroadcasts to\nmobile clients.", ACCENT_BLUE),
    ("TIER 3", "MOBILE APP", "Flutter Ground\nControl Interface\n\nLive video, AI overlay,\nGPS tracking,\ntelemetry, joystick\ncontrol, mission\nmanagement.", ACCENT_PURPLE),
]

for i, (badge, title, desc, accent) in enumerate(tiers):
    left = Inches(0.8 + i * 4.2)
    top = Inches(1.8)
    
    # Card
    card = add_shape(slide, left, top, Inches(3.8), Inches(5.2), BG_CARD)
    
    # Top accent
    add_shape(slide, left, top, Inches(3.8), Inches(0.06), accent)
    
    # Badge
    badge_shape = add_shape(slide, left + Inches(0.3), top + Inches(0.3), Inches(1.2), Inches(0.35), accent)
    tf = badge_shape.text_frame
    p = tf.paragraphs[0]
    p.text = badge
    p.font.size = Pt(9)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_textbox(slide, left + Inches(0.3), top + Inches(0.8), Inches(3.2), Inches(0.4),
                title, font_size=20, color=WHITE, bold=True)
    
    # Divider
    add_shape(slide, left + Inches(0.3), top + Inches(1.3), Inches(3.2), Inches(0.02), accent)
    
    add_textbox(slide, left + Inches(0.3), top + Inches(1.5), Inches(3.2), Inches(3.5),
                desc, font_size=13, color=LIGHT_GRAY)

# Arrows between tiers - simple rectangles
for i in range(2):
    left = Inches(4.6 + i * 4.2)
    add_textbox(slide, left, Inches(3.8), Inches(0.4), Inches(0.5),
                "▶", font_size=24, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Data flow label
add_textbox(slide, Inches(2.5), Inches(7.0), Inches(8), Inches(0.4),
            "Drone → TCP Video → Laptop Server → WebSocket (Video + Detections) → Mobile App",
            font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_page_number(slide, 5, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 6: HARDWARE COMPONENTS
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Hardware Components", "Raspberry Pi Onboard Drone Module")

hw_items = [
    ("Raspberry Pi", "Onboard computer processing\ndrone data & communications", "Pi 3/4"),
    ("Camera Module", "Live video capture for\naerial surveillance feed", "Pi Camera v2"),
    ("GPS Module", "Real-time positioning data\n(NEO-6M) via UART", "NEO-6M"),
    ("Drone Frame", "Quadcopter platform with\npayload capacity for\nRPi + camera + GPS", "Custom / DJI"),
]

for i, (name, desc, spec) in enumerate(hw_items):
    left = Inches(0.8 + i * 3.15)
    top = Inches(1.8)
    card = add_shape(slide, left, top, Inches(2.8), Inches(3.0), BG_CARD)
    # Colored header
    colors = [ACCENT_CYAN, ACCENT_BLUE, ACCENT_PURPLE, GREEN]
    add_shape(slide, left, top, Inches(2.8), Inches(0.04), colors[i])
    
    # Icon placeholder
    icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
        left + Inches(1.0), top + Inches(0.3), Inches(0.8), Inches(0.8))
    icon.fill.solid()
    icon.fill.fore_color.rgb = colors[i]
    icon.line.fill.background()
    tf = icon.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(24)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_textbox(slide, left + Inches(0.2), top + Inches(1.3), Inches(2.4), Inches(0.3),
                name, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), top + Inches(1.7), Inches(2.4), Inches(0.8),
                desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.2), top + Inches(2.5), Inches(2.4), Inches(0.3),
                f"({spec})", font_size=10, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Data flow diagram at bottom
add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.04), DARK_GRAY)
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.5),
            "Data Flow:  GPS Module → Raspberry Pi → TCP Socket → Laptop Server\n"
            "             Camera → Raspberry Pi → TCP Video Stream → FastAPI Server",
            font_size=14, color=DARK_GRAY)

add_page_number(slide, 6, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 7: BACKEND SERVER (FastAPI)
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Backend Server", "FastAPI Middleware & AI Processing Pipeline")

# Pipeline steps
steps = [
    ("TCP Receiver\n(Port 9000)", "Receives raw\nvideo frames\nfrom drone\nvia TCP\nsocket", ACCENT_CYAN),
    ("Frame\nQueue", "Buffers frames\nfor processing\nat controlled\nrate", ACCENT_BLUE),
    ("AI Pipeline\n(YOLO)", "Resize, detect,\n& encode frames\nwith bounding\nboxes", ACCENT_PURPLE),
    ("Output\nQueue", "Queues\nprocessed frames\nfor broadcast\nto clients", GREEN),
    ("WebSocket\nBroadcaster", "Streams video &\ndetections to\nall connected\nmobile clients", ORANGE),
]

for i, (title, desc, color) in enumerate(steps):
    left = Inches(0.5 + i * 2.55)
    top = Inches(2.0)
    
    # Step box
    card = add_shape(slide, left, top, Inches(2.2), Inches(2.8), BG_CARD)
    add_shape(slide, left, top, Inches(2.2), Inches(0.04), color)
    
    add_textbox(slide, left + Inches(0.15), top + Inches(0.2), Inches(1.9), Inches(0.8),
                title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_shape(slide, left + Inches(0.3), top + Inches(1.0), Inches(1.6), Inches(0.015), DARK_GRAY)
    add_textbox(slide, left + Inches(0.15), top + Inches(1.2), Inches(1.9), Inches(1.4),
                desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    # Arrow
    if i < len(steps) - 1:
        add_textbox(slide, left + Inches(2.1), top + Inches(1.0), Inches(0.4), Inches(0.4),
                    "→", font_size=18, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# WebSocket channels info
channels = [
    ("/ws/video", "Binary (JPEG)", "Processed video frames at ~15-30 FPS"),
    ("/ws/detections", "JSON", "AI detection results: bounding boxes, confidence, thermal"),
    ("/ws/commands", "JSON (bidirectional)", "Drone control commands: move, mission, land"),
]
for i, (ep, fmt, desc) in enumerate(channels):
    add_textbox(slide, Inches(1.0 + i * 4.0), Inches(5.3), Inches(3.5), Inches(0.3),
                ep, font_size=14, color=ACCENT_CYAN, bold=True)
    add_textbox(slide, Inches(1.0 + i * 4.0), Inches(5.6), Inches(3.5), Inches(0.3),
                fmt, font_size=11, color=DARK_GRAY)
    add_textbox(slide, Inches(1.0 + i * 4.0), Inches(5.9), Inches(3.5), Inches(0.6),
                desc, font_size=11, color=LIGHT_GRAY)

add_page_number(slide, 7, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 8: AI DETECTION
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "AI Detection Engine", "YOLO-Based Human Detection Pipeline")

# Left: How it works
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(5.5), Inches(0.4),
            "Detection Pipeline", font_size=22, color=WHITE, bold=True)

steps_ai = [
    "1.  Frame Capture — Extract frame from video stream",
    "2.  Preprocessing — Resize to 640×640, normalize pixels",
    "3.  YOLO Inference — Forward pass through CNN model",
    "4.  Post-processing — NMS filtering, confidence thresholding",
    "5.  Bounding Boxes — Draw detections with labels & scores",
    "6.  Encoding — JPEG encode for WebSocket transmission",
]
add_bullet_textbox(slide, Inches(1.0), Inches(2.3), Inches(5.5), Inches(3.0),
                   steps_ai, font_size=14, color=LIGHT_GRAY, spacing=Pt(8))

# Right: Features card
card = add_shape(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), BG_CARD)
add_shape(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.04), ACCENT_PURPLE)
add_textbox(slide, Inches(7.3), Inches(2.1), Inches(5.0), Inches(0.4),
            "Detection Features", font_size=20, color=ACCENT_PURPLE, bold=True)

features_ai = [
    "• Real-time inference at ~15-30 FPS",
    "• Confidence score for each detection",
    "• Thermal mode simulation via color filters",
    "• HUD overlay with bounding boxes & labels",
    "• Human count tracking per frame",
    "• Auto-generated incident reports on detection",
    "• Fallback dummy detector for testing",
]
add_bullet_textbox(slide, Inches(7.3), Inches(2.7), Inches(5.0), Inches(3.5),
                   features_ai, font_size=14, color=LIGHT_GRAY, spacing=Pt(6))

add_page_number(slide, 8, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 9: MOBILE APP OVERVIEW
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Mobile Application", "Flutter Ground Control Interface")

app_features = [
    ("Live Video", "Real-time MJPEG\nstreaming with\nNormal, Thermal &\nAI Overlay modes", ACCENT_CYAN),
    ("GPS Tracking", "OpenStreetMap with\ndrone position,\npath history &\nuser location", ACCENT_BLUE),
    ("Telemetry", "Circular indicators\nfor battery, altitude,\nspeed, temperature\n& human count", ACCENT_PURPLE),
    ("Drone Control", "Full-screen virtual\njoystick for movement,\nmission start/stop\n& landing commands", GREEN),
    ("Reports", "Real-time incident\nreports: human\ndetected, overheat,\nmission complete", ORANGE),
    ("Authentication", "Email/Password,\nGoogle OAuth &\nOTP verification\nwith Firebase", RED),
]

for i, (title, desc, color) in enumerate(app_features):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.0)
    top = Inches(1.9 + row * 2.6)
    add_icon_card(slide, left, top, Inches(3.6), Inches(2.2),
                  "", title, desc, color)

add_page_number(slide, 9, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 10: AUTHENTICATION SYSTEM
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Authentication System", "Secure User Access with Firebase Auth")

# Left side - auth methods
add_textbox(slide, Inches(1.0), Inches(1.9), Inches(5), Inches(0.4),
            "Authentication Methods", font_size=22, color=WHITE, bold=True)

auth_methods = [
    "📧  Email & Password Registration / Login",
    "    — Secure hashed credentials via Firebase Auth",
    "",
    "🌐  Google OAuth Integration",
    "    — One-tap sign-in with existing Google account",
    "",
    "🍎  Apple Sign-In (Placeholder)",
    "    — Cross-platform support for iOS users",
    "",
    "🔑  OTP Verification",
    "    — 4-digit code sent to email for verification",
    "",
    "🔄  Password Reset Flow",
    "    — Forgot password → Email link → New password",
]
add_bullet_textbox(slide, Inches(1.0), Inches(2.4), Inches(5.5), Inches(4.5),
                   auth_methods, font_size=13, color=LIGHT_GRAY, spacing=Pt(2))

# Right side - flow card
card = add_shape(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(5.0), BG_CARD)
add_shape(slide, Inches(7.2), Inches(1.9), Inches(5.3), Inches(0.04), ACCENT_CYAN)
add_textbox(slide, Inches(7.5), Inches(2.2), Inches(4.8), Inches(0.4),
            "Screens Flow", font_size=20, color=WHITE, bold=True)

flow_items = [
    "→ Splash Screen (auto-auth check)",
    "→ Onboarding (4-page carousel)",
    "→ Sign Up / Sign In",
    "→ OTP Verification",
    "→ Forgot / New Password",
    "→ Home (Drone Dashboard)",
]
add_bullet_textbox(slide, Inches(7.5), Inches(2.8), Inches(4.8), Inches(3.5),
                   flow_items, font_size=15, color=LIGHT_GRAY, spacing=Pt(8))

# Data row
add_shape(slide, Inches(7.2), Inches(5.5), Inches(5.3), Inches(0.02), DARK_GRAY)
add_textbox(slide, Inches(7.5), Inches(5.7), Inches(4.8), Inches(1.0),
            "Backend: Firebase Authentication + Cloud Firestore\n"
            "State: AuthCubit (BLoC pattern) with sealed states\n"
            "DI: GetIt service locator for auth repository",
            font_size=12, color=DARK_GRAY)

add_page_number(slide, 10, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 11: ONBOARDING & SPLASH
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Onboarding & Splash", "First-Time User Experience")

# Splash card
card = add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5), BG_CARD)
add_shape(slide, Inches(0.8), Inches(1.8), Inches(0.04), Inches(2.5), ACCENT_CYAN)
add_textbox(slide, Inches(1.2), Inches(1.9), Inches(4.8), Inches(0.3),
            "Animated Splash Screen", font_size=18, color=WHITE, bold=True)
add_textbox(slide, Inches(1.2), Inches(2.3), Inches(4.8), Inches(1.8),
            "• Phoenix logo animation with fade-in + scale\n"
            "• SVG-based logo rendering\n"
            "• Automatic auth check on startup\n"
            "• Routes to Onboarding or Home accordingly",
            font_size=13, color=LIGHT_GRAY)

# Onboarding card
card = add_shape(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(2.5), BG_CARD)
add_shape(slide, Inches(0.8), Inches(4.6), Inches(0.04), Inches(2.5), ACCENT_BLUE)
add_textbox(slide, Inches(1.2), Inches(4.7), Inches(4.8), Inches(0.3),
            "4-Page Onboarding Carousel", font_size=18, color=WHITE, bold=True)
add_textbox(slide, Inches(1.2), Inches(5.1), Inches(4.8), Inches(1.8),
            "• Page 1: Saving Lives — Mission statement\n"
            "• Page 2: Thermal Scanning — AI detection\n"
            "• Page 3: Stay Connected — Real-time data\n"
            "• Page 4: Every Second Counts — Call to action\n"
            "• Interactive PageView with dot indicators",
            font_size=13, color=LIGHT_GRAY)

# Right side - tech details
card = add_shape(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(2.0), BG_CARD)
add_shape(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.04), ACCENT_PURPLE)
add_textbox(slide, Inches(7.3), Inches(2.0), Inches(5.0), Inches(0.3),
            "Implementation Details", font_size=16, color=WHITE, bold=True)
add_textbox(slide, Inches(7.3), Inches(2.4), Inches(5.0), Inches(1.2),
            "• BlocSelector for efficient rebuilds\n"
            "• Skip/Next/Start buttons with state management\n"
            "• SVG assets stored as embedded base64\n"
            "• Page snap physics for natural feel",
            font_size=12, color=LIGHT_GRAY)

# Screenshots placeholder
card = add_shape(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(2.9), BG_CARD)
add_shape(slide, Inches(7.0), Inches(4.2), Inches(5.5), Inches(0.04), GREEN)
add_textbox(slide, Inches(7.3), Inches(4.4), Inches(5.0), Inches(0.3),
            "Key Metrics", font_size=16, color=WHITE, bold=True)
metrics = [
    "• Total screens: 10 (auth + onboarding + drone)",
    "• Routing: GoRouter with declarative config",
    "• State management: BLoC/Cubit pattern",
    "• Screen adaptation: flutter_screenutil",
    "• Theme: Professional dark-first design",
]
add_bullet_textbox(slide, Inches(7.3), Inches(4.8), Inches(5.0), Inches(2.0),
                   metrics, font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

add_page_number(slide, 11, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 12: LIVE VIDEO STREAMING
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Live Video Streaming", "MJPEG over WebSocket with Multi-Mode Display")

# Architecture flow
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(0.4),
            "Video Pipeline", font_size=20, color=WHITE, bold=True)

# Flow boxes
flow_steps = [
    "Drone\nCamera", "TCP\nSocket", "FastAPI\nServer", "WebSocket\nBroadcast", "Flutter\nApp"
]
for i, step in enumerate(flow_steps):
    x = Inches(0.8 + i * 2.5)
    box = add_shape(slide, x, Inches(2.4), Inches(2.0), Inches(0.9), BG_CARD)
    add_shape(slide, x, Inches(2.4), Inches(2.0), Inches(0.03), ACCENT_CYAN)
    add_textbox(slide, x, Inches(2.5), Inches(2.0), Inches(0.7),
                step, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        add_textbox(slide, x + Inches(1.95), Inches(2.6), Inches(0.5), Inches(0.4),
                    "→", font_size=18, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

# Three video modes
modes = [
    ("NORMAL MODE", "Standard RGB video\nfeed from drone\ncamera with no\nimage processing", ACCENT_CYAN),
    ("THERMAL MODE", "Simulated thermal\nimagery using 5×5\ncolor filter matrix\nfor heat mapping", ORANGE),
    ("AI OVERLAY", "YOLO detection\nbounding boxes +\nconfidence scores\noverlaid on video", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(modes):
    left = Inches(0.8 + i * 4.2)
    top = Inches(3.8)
    card = add_shape(slide, left, top, Inches(3.8), Inches(2.5), BG_CARD)
    add_shape(slide, left, top, Inches(3.8), Inches(0.04), color)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.2), Inches(3.2), Inches(0.3),
                title, font_size=14, color=color, bold=True)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.7), Inches(3.2), Inches(1.5),
                desc, font_size=13, color=LIGHT_GRAY)

# Bottom info
add_shape(slide, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.02), DARK_GRAY)
add_textbox(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
            "WebSocket Client: 3 concurrent channels (video/detection/command)  |  DroneWsBridge converts raw streams to domain events",
            font_size=12, color=DARK_GRAY)

add_page_number(slide, 12, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 13: AI OVERLAY & HUD
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "AI Overlay & HUD", "Real-Time Detection Visualization")

# Left side - overlay features
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(6), Inches(0.4),
            "On-Screen Overlay Elements", font_size=20, color=WHITE, bold=True)

overlay_items = [
    "• Bounding Boxes — Colored rectangles around detected persons",
    "• Confidence Scores — Percentage label above each box",
    "• Human Count — Live counter in top corner",
    "• Grid Pattern — Custom GridPainter overlay for spatial reference",
    "• Gradient Overlay — Transparent gradient for readability",
    "• Status Bar — Connection status + telemetry strip",
]
add_bullet_textbox(slide, Inches(1.0), Inches(2.3), Inches(6), Inches(3.0),
                   overlay_items, font_size=14, color=LIGHT_GRAY, spacing=Pt(8))

# Right side - technical card
card = add_shape(slide, Inches(7.5), Inches(1.8), Inches(5.0), Inches(5.0), BG_CARD)
add_shape(slide, Inches(7.5), Inches(1.8), Inches(5.0), Inches(0.04), ACCENT_PURPLE)
add_textbox(slide, Inches(7.8), Inches(2.1), Inches(4.5), Inches(0.4),
            "AI Overlay Technical Stack", font_size=18, color=ACCENT_PURPLE, bold=True)

tech_overlay = [
    "Detection Data (JSON from WebSocket):",
    "",
    "  {",
    '    "detections": [',
    "      {",
    '        "label": "person",',
    '        "confidence": 0.92,',
    '        "x": 0.45, "y": 0.32,',
    '        "w": 0.12, "h": 0.28',
    "      }",
    "    ],",
    '    "human_count": 3',
    "  }",
    "",
    "Rendering: Custom AiDetectionOverlay widget",
    "with CustomPainter for bounding boxes",
]
add_bullet_textbox(slide, Inches(7.8), Inches(2.6), Inches(4.5), Inches(4.0),
                   tech_overlay, font_size=10, color=LIGHT_GRAY, spacing=Pt(2))

add_page_number(slide, 13, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 14: GPS TRACKING & MAPS
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "GPS Tracking & Mapping", "Real-Time Drone Position on OpenStreetMap")

# Left - map features
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(6), Inches(0.4),
            "Map Features", font_size=20, color=WHITE, bold=True)

map_features = [
    "• OpenStreetMap Integration via flutter_map",
    "• Real-time drone marker with red pin",
    "• User location marker with blue dot",
    "• Path history — polyline track of drone movement",
    "• Auto-follow: camera centers on drone position",
    "• Responsive tile loading with caching",
]
add_bullet_textbox(slide, Inches(1.0), Inches(2.3), Inches(5.5), Inches(3.0),
                   map_features, font_size=14, color=LIGHT_GRAY, spacing=Pt(8))

# Right - tech card
card = add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(2.3), BG_CARD)
add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(0.04), ACCENT_CYAN)
add_textbox(slide, Inches(7.5), Inches(2.0), Inches(4.8), Inches(0.3),
            "Data Source", font_size=16, color=WHITE, bold=True)
add_textbox(slide, Inches(7.5), Inches(2.4), Inches(4.8), Inches(1.5),
            "GPS Coordinates published to Firestore:\n"
            "  drone/location → {lat, lng, timestamp}\n\n"
            "DroneTrackingCubit subscribes to real-time\n"
            "Firestore stream with 10s disconnect timer",
            font_size=13, color=LIGHT_GRAY)

# Bottom - status indicators
card = add_shape(slide, Inches(7.2), Inches(4.4), Inches(5.3), Inches(2.4), BG_CARD)
add_shape(slide, Inches(7.2), Inches(4.4), Inches(5.3), Inches(0.04), GREEN)
add_textbox(slide, Inches(7.5), Inches(4.6), Inches(4.8), Inches(0.3),
            "Map Screen Components", font_size=16, color=WHITE, bold=True)
map_components = [
    "• MapScreenHeader — Greeting + connection status",
    "• DroneMapView — Map controller + state handling",
    "• DroneMapStatusBar — Battery/Temp/Signal row",
    "• MissionStatusCard — Pre-completed missions count",
    "• StartMissionButton — Gradient action button",
]
add_bullet_textbox(slide, Inches(7.5), Inches(5.0), Inches(4.8), Inches(1.5),
                   map_components, font_size=12, color=LIGHT_GRAY, spacing=Pt(4))

add_page_number(slide, 14, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 15: TELEMETRY DASHBOARD
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Telemetry Monitoring", "Real-Time Drone Health & Status Dashboard")

# Telemetry indicators
telemetry = [
    ("BATTERY", "75%", "Remaining\npower level", ACCENT_CYAN),
    ("HUMANS", "3", "Detected\npersons", ACCENT_BLUE),
    ("ALTITUDE", "42m", "Height above\nground", ACCENT_PURPLE),
    ("SPEED", "8.5 m/s", "Current\ntravel speed", GREEN),
    ("TEMP", "38°C", "System\ntemperature", ORANGE),
]

for i, (label, value, desc, color) in enumerate(telemetry):
    left = Inches(0.5 + i * 2.55)
    top = Inches(1.9)
    
    card = add_shape(slide, left, top, Inches(2.3), Inches(2.8), BG_CARD)
    add_shape(slide, left, top, Inches(2.3), Inches(0.04), color)
    
    # Circular indicator placeholder
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        left + Inches(0.65), top + Inches(0.3), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    # Make semi-transparent by layering
    inner = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        left + Inches(0.75), top + Inches(0.4), Inches(0.8), Inches(0.8))
    inner.fill.solid()
    inner.fill.fore_color.rgb = BG_CARD
    inner.line.fill.background()
    
    add_textbox(slide, left + Inches(0.1), top + Inches(0.5), Inches(2.1), Inches(0.5),
                value, font_size=22, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(1.5), Inches(2.1), Inches(0.3),
                label, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(1.9), Inches(2.1), Inches(0.6),
                desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom section - status card
card = add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(2.0), BG_CARD)
add_shape(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.04), ACCENT_CYAN)
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(5), Inches(0.3),
            "Connection Status & Reports", font_size=16, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.8), Inches(5.5), Inches(1.2),
            "• ConnectionStatusBar — Green/Red dot indicator\n"
            "• DroneStatusCard — Summary card with circular indicators + mini map\n"
            "• ReportsSection — Real-time incident feed with timestamps\n"
            "• Report types: Human Detected, System Overheated, Mission Complete",
            font_size=12, color=LIGHT_GRAY)
add_textbox(slide, Inches(7.5), Inches(5.8), Inches(5), Inches(1.2),
            "Data Source:\n"
            "Firestore real-time stream from drone/status\n"
            "DroneStatusCubit manages telemetry state\n"
            "Automatic disconnection detection (10s timer)",
            font_size=12, color=DARK_GRAY)

add_page_number(slide, 15, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 16: VIRTUAL JOYSTICK
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Virtual Joystick", "Touch-Based Drone Control Interface")

# Left side - joystick features
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(6), Inches(0.4),
            "Joystick Capabilities", font_size=20, color=WHITE, bold=True)

joy_features = [
    "• Full-screen immersive touch interface",
    "• Custom JoystickPainter with gradient thumb",
    "• Directional indicators around the thumb",
    "• Normalized output (x, y) for drone commands",
    "• Tap-to-toggle visibility of controls overlay",
    "• Gesture-based pan with smooth tracking",
]
add_bullet_textbox(slide, Inches(1.0), Inches(2.3), Inches(5.5), Inches(3.0),
                   joy_features, font_size=14, color=LIGHT_GRAY, spacing=Pt(8))

# Right side - command mapping
card = add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(5.0), BG_CARD)
add_shape(slide, Inches(7.2), Inches(1.8), Inches(5.3), Inches(0.04), GREEN)
add_textbox(slide, Inches(7.5), Inches(2.1), Inches(4.8), Inches(0.4),
            "Command Mapping", font_size=18, color=WHITE, bold=True)

commands_info = [
    "WebSocket Command Format:",
    "",
    '  {"cmd": "move", "data": {"x": 0.5, "y": -0.3}}',
    '  {"cmd": "start_mission", "data": {}}',
    '  {"cmd": "stop_mission", "data": {}}',
    '  {"cmd": "land", "data": {}}',
    "",
    "Commands written to:",
    "  Firestore → drone/commands ",
    "  WebSocket → /ws/commands ",
    "",
    "Dual-path ensures reliability:",
    "  • Firestore for persistent commands",
    "  • WebSocket for low-latency control",
]
add_bullet_textbox(slide, Inches(7.5), Inches(2.6), Inches(4.8), Inches(4.0),
                   commands_info, font_size=12, color=LIGHT_GRAY, spacing=Pt(3))

add_page_number(slide, 16, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 17: MISSION MANAGEMENT
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Mission Management", "End-to-End Rescue Mission Lifecycle")

# Mission lifecycle
phases = [
    ("PLAN", "Define mission\nparameters &\nsearch area\ncoordinates", ACCENT_CYAN),
    ("DEPLOY", "Launch drone &\nestablish\nautonomous\nflight path", ACCENT_BLUE),
    ("MONITOR", "Real-time video\n+ AI detection\n+ telemetry\nstreaming", ACCENT_PURPLE),
    ("REPORT", "Automated\nincident reports\n& mission\nsummary logs", ORANGE),
    ("COMPLETE", "Return to base\nmission debrief\n& data archive", GREEN),
]

for i, (phase, desc, color) in enumerate(phases):
    left = Inches(0.5 + i * 2.55)
    top = Inches(1.9)
    
    card = add_shape(slide, left, top, Inches(2.3), Inches(3.0), BG_CARD)
    add_shape(slide, left, top, Inches(2.3), Inches(0.04), color)
    
    # Phase number
    num_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        left + Inches(0.85), top + Inches(0.25), Inches(0.6), Inches(0.6))
    num_shape.fill.solid()
    num_shape.fill.fore_color.rgb = color
    num_shape.line.fill.background()
    tf = num_shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    add_textbox(slide, left + Inches(0.1), top + Inches(1.0), Inches(2.1), Inches(0.3),
                phase, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), top + Inches(1.4), Inches(2.1), Inches(1.2),
                desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    
    if i < len(phases) - 1:
        add_textbox(slide, left + Inches(2.2), top + Inches(1.2), Inches(0.3), Inches(0.3),
                    "→", font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Mission features at bottom
card = add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), BG_CARD)
add_shape(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(0.04), ACCENT_CYAN)
add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.3),
            "Key Features", font_size=16, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0),
            "• PLAN MISSION button on home screen  •  Start Mission on map screen  •  Real-time mission status tracking  •  Automated incident reports\n"
            "• Mission history and pre-completed mission count  •  Stop/Land commands for emergency abort  •  Firebase Firestore for persistent mission data",
            font_size=12, color=LIGHT_GRAY)

add_page_number(slide, 17, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 18: REPORTS & ALERTS
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Reports & Alerts", "Real-Time Incident Notification System")

# Report types
report_types = [
    ("🚨", "Human Detected", "AI detection triggered\nwith confidence score\nand timestamp", RED),
    ("🌡", "System Overheated", "Temperature threshold\nexceeded warning\nnotification", ORANGE),
    ("✅", "Mission Complete", "Successful mission\ncompletion summary\nreport", GREEN),
    ("⚠", "Connection Lost", "Drone disconnection\nalert with automatic\nreconnect timer", ACCENT_CYAN),
]

for i, (icon, title, desc, color) in enumerate(report_types):
    left = Inches(0.5 + i * 3.2)
    top = Inches(1.9)
    
    card = add_shape(slide, left, top, Inches(2.9), Inches(2.5), BG_CARD)
    add_shape(slide, left, top, Inches(2.9), Inches(0.04), color)
    
    add_textbox(slide, left + Inches(0.2), top + Inches(0.3), Inches(2.5), Inches(0.4),
                icon, font_size=28, color=color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.8), Inches(2.5), Inches(0.3),
                title, font_size=16, color=color, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(1.2), Inches(2.5), Inches(1.0),
                desc, font_size=11, color=LIGHT_GRAY)

# Bottom - report system details
card = add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.4), BG_CARD)
add_shape(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.04), ACCENT_PURPLE)
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(5), Inches(0.3),
            "Report System Architecture", font_size=16, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(5.4), Inches(5.5), Inches(1.5),
            "• Reports stored in Firestore subcollection:\n"
            "  drone/reports/entries/{document}\n"
            "• Each report has: type, message, timestamp\n"
            "• ReportsSection widget with empty state\n"
            "• ReportTile with timeAgo formatting\n"
            "• Auto-generated from AI detections & system events",
            font_size=12, color=LIGHT_GRAY)

add_textbox(slide, Inches(7.5), Inches(5.4), Inches(5), Inches(1.5),
            "Report Types Implementation:\n\n"
            "  enum ReportType {\n"
            "    humanDetected,\n"
            "    systemOverheated,\n"
            "    missionComplete,\n"
            "    connectionLost\n"
            "  }",
            font_size=11, color=LIGHT_GRAY)

add_page_number(slide, 18, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 19: UI & DESIGN HIGHLIGHTS
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "UI/UX Design Highlights", "Professional Dark Interface for Field Operations")

# Design aspects
design_aspects = [
    ("Dark-First Theme", "Optimized for outdoor\nfield use with reduced\nblue light and high\ncontrast readability", ACCENT_CYAN),
    ("Responsive Layout", "flutter_screenutil\nadapts to phones\nand tablets (base:\n393×852)", ACCENT_BLUE),
    ("Smooth Animations", "Animated transitions\nfor snackbars, page\nnavigation, state\nchanges", ACCENT_PURPLE),
    ("SVG Assets", "Custom SVG icons and\nlogo with embedded\nbase64 rendering\nvia flutter_svg", GREEN),
    ("Centralized Theme", "AppTheme + AppDimensions\nfor consistent colors,\ngradients, spacing,\nand typography", ORANGE),
    ("Custom Painters", "GridPainter for\nvideo overlay,\nJoystickPainter\nfor controls", RED),
]

for i, (title, desc, color) in enumerate(design_aspects):
    row = i // 3
    col = i % 3
    left = Inches(0.8 + col * 4.0)
    top = Inches(1.9 + row * 2.6)
    add_icon_card(slide, left, top, Inches(3.6), Inches(2.2),
                  "", title, desc, color)

add_page_number(slide, 19, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 20: STATE MANAGEMENT & ARCHITECTURE
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Architecture & State Management", "Clean Architecture with BLoC Pattern")

# Architecture layers
layers = [
    ("DATA LAYER", "Repository Implementations\nFirebase Auth & Firestore\nWebSocket Client\nDio HTTP Client\nSharedPreferences", ACCENT_CYAN),
    ("DOMAIN LAYER", "Abstract Repositories\nEntity Models (Equatable)\nFailure Classes\neither<Failure, T>\npatterns", ACCENT_BLUE),
    ("PRESENTATION", "Cubits / BLoCs\nScreens & Widgets\nCustom Painters\nState: sealed classes\nwith copyWith", ACCENT_PURPLE),
]

for i, (title, desc, color) in enumerate(layers):
    left = Inches(0.8 + i * 4.2)
    top = Inches(1.9)
    
    card = add_shape(slide, left, top, Inches(3.8), Inches(3.5), BG_CARD)
    add_shape(slide, left, top, Inches(3.8), Inches(0.06), color)
    
    add_textbox(slide, left + Inches(0.3), top + Inches(0.3), Inches(3.2), Inches(0.3),
                title, font_size=16, color=color, bold=True)
    add_shape(slide, left + Inches(0.3), top + Inches(0.7), Inches(3.2), Inches(0.02), color)
    add_textbox(slide, left + Inches(0.3), top + Inches(0.9), Inches(3.2), Inches(2.3),
                desc, font_size=13, color=LIGHT_GRAY)

# Bottom - Cubit listing
card = add_shape(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.5), BG_CARD)
add_shape(slide, Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.04), ACCENT_CYAN)
add_textbox(slide, Inches(1.1), Inches(5.9), Inches(11.2), Inches(0.3),
            "Cubits:  OnboardingCubit | AuthCubit | DroneTrackingCubit | DroneStatusCubit | VideoFeedCubit",
            font_size=14, color=WHITE, bold=True)
add_textbox(slide, Inches(1.1), Inches(6.3), Inches(11.2), Inches(0.7),
            "DI: GetIt service locator  •  Routing: GoRouter (10 named routes)  •  Error handling: dartz Either type  •  ~70+ Dart files across the project",
            font_size=12, color=DARK_GRAY)

add_page_number(slide, 20, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 21: TECH STACK
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Technology Stack", "Tools & Technologies Used")

# Tech categories
categories = [
    ("MOBILE APP", [
        ("Flutter 3.12+", "Cross-platform UI framework"),
        ("Dart 3.12+", "Programming language"),
        ("BLoC/Cubit", "State management"),
        ("GoRouter", "Declarative routing"),
        ("get_it", "Dependency injection"),
        ("flutter_map", "OpenStreetMap integration"),
        ("web_socket_channel", "Real-time communication"),
        ("flutter_screenutil", "Responsive design"),
    ], ACCENT_CYAN),
    ("BACKEND", [
        ("FastAPI", "Python web framework"),
        ("Uvicorn", "ASGI server"),
        ("OpenCV", "Video processing"),
        ("NumPy / Pillow", "Image manipulation"),
        ("YOLO", "AI object detection"),
        ("WebSockets", "Real-time broadcast"),
        ("Firebase Admin", "Server-side Firebase"),
    ], ACCENT_BLUE),
    ("INFRASTRUCTURE", [
        ("Firebase Auth", "Authentication"),
        ("Cloud Firestore", "Real-time database"),
        ("Firebase Core", "Platform initialization"),
        ("SharedPreferences", "Local storage"),
        ("Geolocator", "Location services"),
        ("google_sign_in", "Google OAuth"),
        ("dartz", "Functional error handling"),
    ], ACCENT_PURPLE),
    ("HARDWARE", [
        ("Raspberry Pi", "Onboard computer"),
        ("Pi Camera v2", "Video capture"),
        ("NEO-6M GPS", "Position tracking"),
        ("TCP Socket", "Video transmission"),
        ("Pi GPIO/UART", "Hardware interface"),
    ], GREEN),
]

for i, (title, items, color) in enumerate(categories):
    left = Inches(0.4 + i * 3.2)
    top = Inches(1.7)
    
    card = add_shape(slide, left, top, Inches(3.0), Inches(5.5), BG_CARD)
    add_shape(slide, left, top, Inches(3.0), Inches(0.04), color)
    
    add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.6), Inches(0.3),
                title, font_size=14, color=color, bold=True)
    add_shape(slide, left + Inches(0.2), top + Inches(0.55), Inches(2.6), Inches(0.02), color)
    
    y_offset = Inches(0.7)
    for name, desc in items:
        add_textbox(slide, left + Inches(0.2), top + y_offset, Inches(2.6), Inches(0.25),
                    f"• {name}", font_size=11, color=WHITE, bold=True)
        add_textbox(slide, left + Inches(0.4), top + y_offset + Inches(0.2), Inches(2.4), Inches(0.25),
                    desc, font_size=9, color=DARK_GRAY)
        y_offset += Inches(0.55)

add_page_number(slide, 21, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 22: CHALLENGES & SOLUTIONS
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Challenges & Solutions", "Key Technical Hurdles Overcome")

challenges = [
    ("Real-Time Video Latency", "Optimized MJPEG streaming\nover WebSocket with frame\nqueue & rate limiting",
     "TCP → Frame Queue → AI Pipeline → WS Broadcast\nachieved ~15-30 FPS with <500ms latency", ACCENT_CYAN),
    ("Disconnection Handling", "10-second grace timer before\nmarking drone as disconnected\nprevents UI flickering",
     "DroneTrackingCubit timer resets on each\nincoming location update from Firestore", ACCENT_BLUE),
    ("Cross-Platform UI Consistency", "flutter_screenutil for adaptive\nlayout across phone & tablet\nwith dark-first design",
     "Base design size: 393×852. All dimensions\nuse .w, .h, .r, .sp extensions", ACCENT_PURPLE),
    ("AI Detection Integration", "Custom pipeline with fallback\ndummy detector for testing\nwithout real drone hardware",
     "Abstract detection interface allows swapping\nYOLO model or using simulated data", GREEN),
    ("Firebase Real-Time Sync", "Multiple Firestore streams for\nlocation, status, reports &\ncommands with error handling",
     "Each stream independently error-handled\nwith Either pattern for clean state", ORANGE),
]

for i, (problem, solution, result, color) in enumerate(challenges):
    top = Inches(1.8 + i * 1.1)
    
    # Number
    num = slide.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(0.5), top + Inches(0.15), Inches(0.4), Inches(0.4))
    num.fill.solid()
    num.fill.fore_color.rgb = color
    num.line.fill.background()
    tf = num.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.color.rgb = BG_DARK
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Problem
    add_textbox(slide, Inches(1.1), top, Inches(3.5), Inches(0.8),
                problem, font_size=14, color=WHITE, bold=True)
    
    # Solution
    add_textbox(slide, Inches(4.8), top, Inches(4.0), Inches(0.8),
                solution, font_size=11, color=LIGHT_GRAY)
    
    # Result
    add_textbox(slide, Inches(9.0), top, Inches(4.0), Inches(0.8),
                result, font_size=11, color=DARK_GRAY)
    
    # Separator line
    if i < len(challenges) - 1:
        add_shape(slide, Inches(0.5), top + Inches(0.95), Inches(12.3), Inches(0.01), DARK_GRAY)

add_page_number(slide, 22, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 23: FUTURE ENHANCEMENTS
# ═══════════════════════════════════════════════
slide = new_slide()
add_section_title(slide, "Future Enhancements", "Roadmap & Next Steps")

enhancements = [
    ("Real YOLO Model Training", "Train custom YOLO model on\nSAR-specific dataset for higher\naccuracy in rescue scenarios", ACCENT_CYAN),
    ("Raspberry Pi Hardware Scripts", "Complete onboard firmware for\nautonomous drone control &\ndata relay to laptop server", ACCENT_BLUE),
    ("Multi-Drone Support", "Simultaneous operation of\nmultiple drones with\nindividual control interfaces", ACCENT_PURPLE),
    ("Offline Mode", "Local caching of video &\ntelemetry data for operation\nin remote areas without internet", GREEN),
    ("Advanced Mission Planning", "GPS waypoint plotting &\nautonomous path planning\nwith obstacle avoidance", ORANGE),
    ("Push Notifications", "Alert first responders when\nhuman detection is triggered\nor system warnings occur", RED),
    ("Flight Recording & Replay", "Record missions for post-flight\nanalysis and training\npurposes", ACCENT_CYAN),
    ("Web Dashboard", "Browser-based monitoring\ninterface for command\ncenter operations", ACCENT_BLUE),
]

for i, (title, desc, color) in enumerate(enhancements):
    row = i // 4
    col = i % 4
    left = Inches(0.4 + col * 3.2)
    top = Inches(1.7 + row * 2.8)
    
    card = add_shape(slide, left, top, Inches(2.9), Inches(2.4), BG_CARD)
    add_shape(slide, left, top, Inches(2.9), Inches(0.04), color)
    
    add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(2.5), Inches(0.4),
                title, font_size=14, color=color, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.7), Inches(2.5), Inches(1.5),
                desc, font_size=11, color=LIGHT_GRAY)

add_page_number(slide, 23, TOTAL_SLIDES)

# ═══════════════════════════════════════════════
# SLIDE 24: THANK YOU
# ═══════════════════════════════════════════════
slide = new_slide()
# Background
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5), BG_DARK)

# Decorative elements
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(5), Inches(5), Inches(5))
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(0x00, 0x65, 0xFF)
circle.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(6), Inches(6))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(0x7C, 0x3A, 0xED)
circle2.line.fill.background()

# Main content
add_textbox(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(1.0),
            "Thank You", font_size=64, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_gradient_line(slide, Inches(5.5), Inches(3.0), Inches(2.3))
add_textbox(slide, Inches(1.0), Inches(3.3), Inches(11.3), Inches(0.6),
            "Questions & Discussion", font_size=28, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.0), Inches(4.5), Inches(11.3), Inches(1.0),
            "Phoenix — Drone Search & Rescue System\n"
            "Project ResQer",
            font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Contact info
add_shape(slide, Inches(0), Inches(6.8), Inches(13.333), Inches(0.02), ACCENT_CYAN)
add_textbox(slide, Inches(1.0), Inches(6.9), Inches(11.3), Inches(0.4),
            "Mahmoud Dahy  |  github.com/mahmoud-dahy  |  Computer Engineering - Graduation Project 2025/2026",
            font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 24, TOTAL_SLIDES)

# ── SAVE ──
output_path = "/home/mahmoud-dahy/Flutter Projects/graduatio_project/Phoenix_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: {os.path.getsize(output_path) / 1024:.1f} KB")
